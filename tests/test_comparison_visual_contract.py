from __future__ import annotations

import contextlib
import hashlib
import importlib.util
import io
import json
import re
import sys
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from tests.test_comparison_report_gate import complete_report as complete_gate_report


ROOT = Path(__file__).resolve().parents[1]


def load_script(name: str, relative_path: str):
    spec = importlib.util.spec_from_file_location(name, ROOT / relative_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Unable to load {relative_path}")
    module = importlib.util.module_from_spec(spec)
    sys.modules[name] = module
    spec.loader.exec_module(module)
    return module


def write_palette_png(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    image = Image.new("P", (1, 1))
    image.putpalette([248, 250, 252] + [0, 0, 0] * 255)
    image.save(path)


def image_mode(path: Path) -> str:
    with Image.open(path) as image:
        return image.mode


class ComparisonVisualContractTests(unittest.TestCase):
    def test_surface_history_breaks_the_second_six_item_catalog_without_losing_content(self) -> None:
        module = load_script("create_mdpr_vs_skill_surface_history", "scripts/create_mdpr_vs_skill_decks.py")
        saturated = ["native-table", "enclosed-card", "enclosed-card", "enclosed-card", "enclosed-card"]
        self.assertEqual(
            module.choose_surface_family("enclosed-card", saturated, alternate="open-catalog"),
            "open-catalog",
        )
        evidence = module.collect_skill_surface_evidence([
            "enclosed-card", "native-table", "enclosed-card", "enclosed-card",
            "enclosed-card", "enclosed-card", "open-catalog", "chart-table", "split-field",
        ])
        self.assertEqual(evidence["maxSameSurfaceRun"], 4)
        self.assertEqual(evidence["saturatedWindows"], [])

        deck = Presentation(ROOT / "artifacts" / "mdpr-vs-skill" / "mdpr-skill-result.pptx")
        first_catalog = deck.slides[4]
        second_catalog = deck.slides[6]
        self.assertEqual(len([shape for shape in first_catalog.shapes if shape.name.startswith("docmap_") and shape.name.endswith("_card")]), 6)
        self.assertEqual(len([shape for shape in second_catalog.shapes if shape.name.startswith("example_") and shape.name.endswith("_card")]), 0)
        expected_titles = ["Basic", "Comparison", "Pipeline", "Diagram", "Five Methods", "Theme Preview"]
        visible_text = [shape.text for shape in second_catalog.shapes if getattr(shape, "has_text_frame", False)]
        for title in expected_titles:
            self.assertEqual(visible_text.count(title), 1)

    def test_runtime_design_evidence_is_copied_and_bounded_from_mdpr_manifest(self) -> None:
        module = load_script("create_mdpr_vs_skill_runtime_evidence", "scripts/create_mdpr_vs_skill_decks.py")
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "build" / "mdpresent-manifest.json"
            stable = root / "artifacts" / "mdpr-runtime-manifest.json"
            source.parent.mkdir(parents=True)
            manifest = {
                "engine": "mdpresent",
                "slideCount": 12,
                "validation": {
                    "polish": {
                        "checked": True,
                        "requiredFailureCount": 0,
                        "chapters": {
                            "layoutComposition": {
                                "required": True,
                                "passed": True,
                                "applicable": True,
                                "eligibleSlideCount": 10,
                                "dominantGeometryRatio": 0.4,
                            },
                        },
                    },
                    "coherence": {
                        "checked": True,
                        "checks": {
                            "claimlessEvidenceSlides": True,
                            "detachedCaptions": True,
                            "orphanTables": True,
                            "lowObjectCoverage": True,
                        },
                        "diagnostics": [{"level": "warning", "code": "EXAMPLE"}],
                    },
                },
            }
            source.write_text(json.dumps(manifest), encoding="utf-8")

            evidence = module.capture_runtime_design_evidence(source, stable, "c" * 40, artifact_root=root)

            self.assertEqual(stable.read_bytes(), source.read_bytes())
            self.assertEqual(evidence["manifestPath"], "artifacts/mdpr-runtime-manifest.json")
            self.assertEqual(evidence["manifestSha256"], hashlib.sha256(stable.read_bytes()).hexdigest())
            self.assertEqual(evidence["polish"]["layoutComposition"]["dominantGeometryRatio"], 0.4)
            self.assertIs(evidence["polish"]["layoutComposition"]["applicable"], True)
            self.assertNotIn("maxSameGeometryInFive", evidence["polish"]["layoutComposition"])
            self.assertEqual(evidence["coherence"]["errorCount"], 0)
            self.assertNotIn("diagnostics", evidence["coherence"])

    def test_generated_body_list_has_no_unassigned_title_band(self) -> None:
        deck = Presentation(ROOT / "artifacts" / "mdpr-vs-skill" / "mdpr-baseline-result.pptx")
        slide = deck.slides[14]
        title_bands = [
            shape
            for shape in slide.shapes
            if not (getattr(shape, "text", "") or "").strip()
            and shape.width >= Inches(9.3)
            and shape.height <= Inches(0.12)
            and Inches(1.3) <= shape.top <= Inches(2.0)
        ]
        text = [
            shape.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text
        ]

        self.assertEqual(title_bands, [])
        for value in [
            "Rendering Rules",
            "- Shared Renderer Contract",
            "- PPTX Renderer",
            "- Decoration Styles",
        ]:
            self.assertEqual(text.count(value), 1)

    def test_mdpr_comparison_table_content_weights_short_label_column(self) -> None:
        deck = Presentation(ROOT / "artifacts" / "mdpr-vs-skill" / "mdpr-baseline-result.pptx")
        slide = deck.slides[3]
        tables = [shape.table for shape in slide.shapes if getattr(shape, "has_table", False)]

        self.assertEqual(len(tables), 1)
        table = tables[0]
        widths = [column.width for column in table.columns]
        total = sum(widths)
        self.assertEqual(len(widths), 3)
        self.assertGreaterEqual(widths[0], Inches(1.35))
        self.assertLessEqual(widths[0], total * 0.24 + 2)
        self.assertLessEqual(abs(widths[1] - widths[2]), 2)

        expected = [
            "Area", "MDPR runtime", "mdpr-skill review companion",
            "Role", "Parser", "Layout", "Output",
        ]
        text = [cell.text for row in table.rows for cell in row.cells]
        for value in expected:
            self.assertEqual(text.count(value), 1)

    def test_sparse_pipeline_continuation_uses_compact_editable_nodes(self) -> None:
        deck = Presentation(ROOT / "artifacts" / "mdpr-vs-skill" / "mdpr-baseline-result.pptx")
        slide = deck.slides[5]
        expected = [
            "MDPR manifest and previews",
            "mdpr-skill hints or review findings",
            "MDPR remains the only renderer",
        ]
        text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
        matched = [next(shape for shape in text_shapes if shape.text == text) for text in expected]
        cards = sorted(
            (
                shape
                for shape in slide.shapes
                if not (getattr(shape, "text", "") or "").strip()
                and Inches(8.5) <= shape.width <= Inches(8.7)
                and shape.height > Inches(0.5)
            ),
            key=lambda shape: shape.top,
        )
        connectors = [
            shape
            for shape in slide.shapes
            if not (getattr(shape, "text", "") or "").strip()
            and shape.width <= Inches(0.01)
            and Inches(0.3) <= shape.height <= Inches(0.5)
        ]

        self.assertEqual([shape.text for shape in matched], expected)
        self.assertEqual(len(cards), 3)
        self.assertEqual([round(shape.height / Inches(1), 2) for shape in cards], [0.95, 0.95, 0.95])
        self.assertAlmostEqual((cards[0].top + cards[-1].top + cards[-1].height) / 2 / Inches(1), 4.195, places=2)
        self.assertEqual(len(connectors), 2)

    def test_paired_comparison_artifacts_preserve_editable_source_ancestry(self) -> None:
        deck = Presentation(ROOT / "artifacts" / "mdpr-vs-skill" / "mdpr-baseline-result.pptx")
        expected = {
            21: [
                ["Current Approach", "Improved Approach"],
                ["Owners manually prepare documents", "Meeting notes and report drafts are generated automatically"],
                ["Quality depends on individual skill", "Documents follow a consistent structure"],
                ["Search and reuse are difficult", "Materials are found with semantic search"],
            ],
            22: [
                ["Current Approach", "Improved Approach"],
                ["Documents are written manually", "Drafts are generated automatically"],
                ["Format varies by person", "Format is standardized by template"],
                ["Search is difficult", "Semantic search is available"],
            ],
        }

        for slide_index, rows in expected.items():
            slide = deck.slides[slide_index]
            tables = [shape.table for shape in slide.shapes if getattr(shape, "has_table", False)]
            self.assertEqual(len(tables), 1)
            actual = [[cell.text for cell in row.cells] for row in tables[0].rows]
            self.assertEqual(actual, rows)

    def test_skill_contract_requires_rendered_visual_revalidation_without_synthetic_rules(self) -> None:
        skill = (ROOT / "skills" / "mdpr-skill" / "SKILL.md").read_text(encoding="utf-8")

        self.assertIn("automatic title underline", skill)
        self.assertIn("isolated bottom rule", skill)
        self.assertIn("Inspect every exported slide", skill)
        self.assertIn("including list\n  and diagram number badges", skill)
        self.assertIn("semantic role and repetition", skill)
        self.assertIn("rendered before/after evidence", skill)
        self.assertIn("text-bearing regions", skill)

    def test_external_mdpr_checkout_path_can_be_recorded_as_evidence(self) -> None:
        module = load_script("create_mdpr_vs_skill_paths", "scripts/create_mdpr_vs_skill_decks.py")
        sibling = ROOT.parent / "mdpresent-spec-scaffold"
        recorded = module.evidence_path(sibling)

        self.assertEqual(recorded, str(sibling.resolve()))

    def test_source_corpus_keeps_manifest_out_of_slide_content_and_normalizes_numbered_titles(self) -> None:
        module = load_script("create_mdpr_vs_skill_corpus", "scripts/create_mdpr_vs_skill_decks.py")
        topic_paths = [
            "docs/01-architecture.md",
            "docs/03-page-splitting.md",
            "docs/04-layout-rules.md",
            "docs/07-rendering-rules.md",
            "docs/11-qa-overflow.md",
        ]
        summaries = [
            {
                "path": path,
                "title": f"{index:02d}. {Path(path).stem}",
                "headingCount": 1,
                "headings": [f"{index:02d}. Topic"],
                "bullets": ["Evidence"],
                "table": [],
                "code": [],
                "codeLanguage": None,
                "charCount": 100,
            }
            for index, path in enumerate(topic_paths, 1)
        ]
        with tempfile.TemporaryDirectory() as temp_dir:
            module.SOURCE_MD = Path(temp_dir) / "corpus.md"
            module.build_source_corpus(summaries)
            corpus = module.SOURCE_MD.read_text(encoding="utf-8")

        self.assertNotIn("## Source manifest", corpus)
        self.assertNotIn("## Parser and splitting topics", corpus)
        self.assertNotIn("## 01.", corpus)
        self.assertIn("## 01-architecture", corpus)

    def test_topic_headings_exclude_normalized_section_title_before_limit(self) -> None:
        module = load_script("create_mdpr_vs_skill_heading_filter", "scripts/create_mdpr_vs_skill_decks.py")
        topic_paths = [
            "docs/01-architecture.md",
            "docs/03-page-splitting.md",
            "docs/04-layout-rules.md",
            "docs/07-rendering-rules.md",
            "docs/11-qa-overflow.md",
        ]
        summaries = []
        for path in topic_paths:
            title = "03. Page Splitting Rules" if path == "docs/03-page-splitting.md" else Path(path).stem
            headings = [title, "Fallback"]
            if path == "docs/03-page-splitting.md":
                headings = [
                    "03 — Page Splitting Rules!",
                    "Heading Rules",
                    "Slide Candidate",
                    "Autosplit Boundary",
                    "AST Parsing",
                    "Overflow",
                ]
            summaries.append({
                "path": path,
                "title": title,
                "headingCount": len(headings),
                "headings": headings,
                "bullets": [],
                "table": [],
                "code": [],
                "codeLanguage": None,
                "charCount": 100,
            })

        with tempfile.TemporaryDirectory() as temp_dir:
            module.SOURCE_MD = Path(temp_dir) / "corpus.md"
            module.build_source_corpus(summaries)
            corpus = module.SOURCE_MD.read_text(encoding="utf-8")

        section = corpus.split("## Page Splitting Rules", 1)[1].split("\n## ", 1)[0]
        self.assertNotIn("- Page Splitting Rules", section)
        retained = [line.removeprefix("- ") for line in section.splitlines() if line.startswith("- ")]
        self.assertEqual(retained, ["Heading Rules", "Slide Candidate", "Autosplit Boundary", "AST Parsing", "Overflow"])

    def test_comparison_source_keeps_heading_ancestry_as_paired_table_columns(self) -> None:
        module = load_script("create_mdpr_vs_skill_hierarchy", "scripts/create_mdpr_vs_skill_decks.py")
        source = "\n".join([
            "# Comparison Structure Example",
            "",
            "## Current Approach and Improved Approach",
            "",
            "### Current Approach",
            "",
            "- Documents are reviewed manually.",
            "- Source location is unclear.",
            "",
            "### Improved Approach",
            "",
            "- Draft from structured source.",
            "- Keep evidence links visible.",
        ])

        outline = module.extract_outline_facts(source)
        groups = module.extract_paired_comparison_groups(outline)

        self.assertEqual(
            [fact["headingPath"] for fact in outline if fact["kind"] == "listItem"],
            [
                ["Comparison Structure Example", "Current Approach and Improved Approach", "Current Approach"],
                ["Comparison Structure Example", "Current Approach and Improved Approach", "Current Approach"],
                ["Comparison Structure Example", "Current Approach and Improved Approach", "Improved Approach"],
                ["Comparison Structure Example", "Current Approach and Improved Approach", "Improved Approach"],
            ],
        )
        self.assertEqual(groups, [
            {"label": "Current Approach", "items": ["Documents are reviewed manually.", "Source location is unclear."]},
            {"label": "Improved Approach", "items": ["Draft from structured source.", "Keep evidence links visible."]},
        ])
        unrelated = module.extract_outline_facts("\n".join([
            "# Workflow Proposal",
            "## Repetitive Work Is Growing",
            "- Meeting cleanup",
            "## Search Costs Are Growing",
            "- Filename search",
        ]))
        self.assertEqual(module.extract_paired_comparison_groups(unrelated), [])
        fenced_example = module.extract_outline_facts("\n".join([
            "# Parser documentation",
            "```markdown",
            "## Current and Improved",
            "### Current Approach",
            "- Manual",
            "### Improved Approach",
            "- Automatic",
            "```",
        ]))
        self.assertEqual(module.extract_paired_comparison_groups(fenced_example), [])

        summaries = [
            {
                "path": path,
                "title": Path(path).stem,
                "headingCount": 1,
                "headings": [Path(path).stem],
                "bullets": [],
                "table": [],
                "code": [],
                "codeLanguage": None,
                "charCount": 100,
            }
            for path in [
                "docs/01-architecture.md",
                "docs/03-page-splitting.md",
                "docs/04-layout-rules.md",
                "docs/07-rendering-rules.md",
                "docs/11-qa-overflow.md",
            ]
        ]
        summaries.append({
            "path": "examples/comparison/deck.md",
            "title": "Comparison Structure Example",
            "headingCount": 4,
            "headings": [fact["text"] for fact in outline if fact["kind"] == "heading"],
            "bullets": [fact["text"] for fact in outline if fact["kind"] == "listItem"],
            "outlineFacts": outline,
            "comparisonGroups": groups,
            "table": [],
            "code": [],
            "codeLanguage": None,
            "charCount": len(source),
        })
        with tempfile.TemporaryDirectory() as temp_dir:
            module.SOURCE_MD = Path(temp_dir) / "corpus.md"
            module.build_source_corpus(summaries)
            corpus = module.SOURCE_MD.read_text(encoding="utf-8")

        section = corpus.split("## Example: examples/comparison/deck.md", 1)[1].split("\n## ", 1)[0]
        self.assertIn("| Current Approach | Improved Approach |", section)
        self.assertIn("| Documents are reviewed manually. | Draft from structured source. |", section)
        self.assertNotIn("- Current Approach\n", section)

    def test_skill_evidence_deck_omits_decorative_one_line_regions_and_uses_readable_type(self) -> None:
        module = load_script("create_mdpr_vs_skill_decks", "scripts/create_mdpr_vs_skill_decks.py")
        summaries = [
            {
                "path": "README.md",
                "title": "MDPR",
                "headingCount": 4,
                "headings": ["MDPR", "Boundary"],
                "bullets": ["Deterministic runtime"],
                "hasTable": False,
                "table": [],
                "hasCode": False,
                "codeLanguage": None,
                "code": [],
                "charCount": 400,
            },
            {
                "path": "docs/01-architecture.md",
                "title": "Architecture",
                "headingCount": 3,
                "headings": ["Architecture", "Pipeline"],
                "bullets": ["Parser", "Layout"],
                "hasTable": False,
                "table": [],
                "hasCode": False,
                "codeLanguage": None,
                "code": [],
                "charCount": 300,
            },
            {
                "path": "examples/basic/deck.md",
                "title": "Example",
                "headingCount": 2,
                "headings": ["Example", "Result"],
                "bullets": ["Editable output"],
                "hasTable": False,
                "table": [],
                "hasCode": False,
                "codeLanguage": None,
                "code": [],
                "charCount": 200,
            },
        ]
        mdpr_result = {"slides": 8, "textFrames": 40, "tables": 1, "charts": 1}

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_root = Path(temp_dir)
            output = temp_root / "comparison.pptx"
            copy = temp_root / "comparison-copy.pptx"
            module.ROOT = temp_root
            module.SOURCE_MD = temp_root / "source.md"
            module.BASELINE_PPTX = temp_root / "baseline.pptx"
            module.SKILL_PPTX = output
            module.SKILL_FROM_MDPR_RUN_PPTX = copy
            module.build_skill_deck(summaries, mdpr_result)

            deck = Presentation(output)
            shape_names = {shape.name for slide in deck.slides for shape in slide.shapes}
            forbidden = {
                "subtitle",
                "title_rule",
                "coverage_band",
                "coverage_note",
                "actual_file_note",
                "icon_caption",
                "chart_note_card",
            }
            self.assertTrue(forbidden.isdisjoint(shape_names), sorted(forbidden & shape_names))
            self.assertGreaterEqual(module.validate_pptx(output)["minFontSizePt"], 16)

            pipeline = deck.slides[3]
            for prefix in ("mdpr", "skill"):
                card = next(shape for shape in pipeline.shapes if shape.name == f"{prefix}_card")
                body = [shape for shape in pipeline.shapes if shape.name.startswith(f"{prefix}_body_")]
                self.assertTrue(body)
                self.assertLessEqual(max(shape.top + shape.height for shape in body), card.top + card.height)

            optional_visual = deck.slides[8]
            optional_names = {shape.name for shape in optional_visual.shapes}
            self.assertNotIn("body_panel", optional_names)
            self.assertNotIn("icon_slot", optional_names)

    def test_preview_reset_removes_stale_render_evidence(self) -> None:
        module = load_script("create_mdpr_vs_skill_preview_reset", "scripts/create_mdpr_vs_skill_decks.py")
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir)
            stale = [
                output / "mdpr_baseline_preview_1.png",
                output / "skill_preview_1.png",
            ]
            for path in stale:
                path.write_bytes(b"stale")

            module.clear_preview_files(output)

            self.assertFalse(any(path.exists() for path in stale))

    def test_actor_role_colors_are_consistent_across_boundary_slides(self) -> None:
        module = load_script("create_mdpr_vs_skill_actor_colors", "scripts/create_mdpr_vs_skill_decks.py")
        deck = Presentation()
        deck.slide_width = module.Inches(module.SLIDE_W)
        deck.slide_height = module.Inches(module.SLIDE_H)
        module.add_actual_mdpr_run_slide(
            deck,
            [{"headingCount": 1, "charCount": 1}],
            {"slides": 1, "textFrames": 1, "tables": 0, "charts": 0},
        )
        module.add_pipeline_slide(deck)
        module.add_text_icon_slide(deck)

        actual = {shape.name: shape for shape in deck.slides[0].shapes}
        pipeline = {shape.name: shape for shape in deck.slides[1].shapes}
        boundary = {shape.name: shape for shape in deck.slides[2].shapes}

        def text_rgb(shape):
            return str(shape.text_frame.paragraphs[0].runs[0].font.color.rgb)

        mdpr_colors = {
            text_rgb(actual["actual_mdpr_title"]),
            text_rgb(pipeline["mdpr_title"]),
            text_rgb(boundary["runtime_heading"]),
            *(text_rgb(boundary[f"decision_{index}_num"]) for index in range(3)),
        }
        skill_colors = {
            text_rgb(actual["actual_skill_title"]),
            text_rgb(pipeline["skill_title"]),
            text_rgb(boundary["suggestion_heading"]),
            *(str(boundary[f"suggestion_{index}_mark"].fill.fore_color.rgb) for index in range(3)),
        }

        self.assertEqual(len(mdpr_colors), 1)
        self.assertEqual(len(skill_colors), 1)
        self.assertNotEqual(mdpr_colors, skill_colors)

    def test_corpus_chart_names_its_docs_and_adr_aggregate(self) -> None:
        module = load_script("create_mdpr_vs_skill_corpus_groups", "scripts/create_mdpr_vs_skill_decks.py")
        summaries = [
            {"path": "docs/guide.md", "headingCount": 2, "charCount": 20},
            {"path": "docs/adr/0001-contract.md", "headingCount": 3, "charCount": 30},
            {"path": "examples/basic/deck.md", "headingCount": 4, "charCount": 40},
            {"path": "README.md", "headingCount": 5, "charCount": 50},
        ]
        deck = Presentation()
        module.add_source_coverage_slide(deck, summaries)
        module.add_chart_slide(deck, summaries)

        coverage = {shape.name: shape for shape in deck.slides[0].shapes}
        self.assertEqual(coverage["group_0_label"].text, "Docs")
        self.assertEqual(coverage["group_0_num"].text, "1")
        self.assertEqual(coverage["group_3_label"].text, "ADR")
        self.assertEqual(coverage["group_3_num"].text, "1")

        chart_slide = deck.slides[1]
        chart = next(shape.chart for shape in chart_slide.shapes if shape.has_chart)
        table = next(shape.table for shape in chart_slide.shapes if shape.has_table)
        self.assertEqual([category.label for category in chart.plots[0].categories], ["Docs + ADR", "Examples", "Root"])
        self.assertEqual([cell.text for cell in table.rows[1].cells], ["Docs + ADR", "2", "5", "50"])

        no_adr = [summary for summary in summaries if "/adr/" not in summary["path"]]
        control = Presentation()
        module.add_chart_slide(control, no_adr)
        control_chart = next(shape.chart for shape in control.slides[0].shapes if shape.has_chart)
        self.assertEqual([category.label for category in control_chart.plots[0].categories][0], "Docs")

    def test_actual_mdpr_neutral_inventory_has_no_comparison_rules_or_blank_band(self) -> None:
        deck = Presentation(ROOT / "artifacts" / "mdpr-vs-skill" / "mdpr-baseline-result.pptx")
        slide = deck.slides[20]
        expected = [
            "basic/deck.md",
            "comparison/deck.md",
            "pipeline/deck.md",
            "diagram-arrangements/deck.md",
            "five-methods/deck.md — Five-Item Layout Example",
            "theme-preview-en/deck.md",
        ]
        self.assertEqual([shape for shape in slide.shapes if getattr(shape, "has_table", False)], [])
        text = "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
        for item in expected:
            self.assertIn(item, text)

        title = next(shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text == "Example decks from MDPR")

        rules = [
            shape for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and (not getattr(shape, "has_text_frame", False) or not shape.text.strip())
            and shape.width >= Inches(2.5)
            and shape.height <= Inches(0.12)
            and shape.top >= title.top + title.height
            and shape.top <= title.top + title.height + Inches(0.8)
        ]
        self.assertEqual(rules, [])

    def test_example_overview_matches_selected_example_families(self) -> None:
        module = load_script("create_mdpr_vs_skill_example_overview", "scripts/create_mdpr_vs_skill_decks.py")
        manifest = json.loads((ROOT / "artifacts" / "mdpr-vs-skill" / "source-manifest.json").read_text(encoding="utf-8"))
        selected_examples = [item for item in manifest["files"] if item["path"].startswith("examples/")]
        selected_families = [
            Path(item["path"]).parent.name.removesuffix("-en")
            for item in selected_examples
        ]
        overview_lines = module.example_overview_lines(selected_examples)
        deck = Presentation(ROOT / "artifacts" / "mdpr-vs-skill" / "mdpr-baseline-result.pptx")
        overview = next(
            slide for slide in deck.slides
            if any(getattr(shape, "text", "") == "Example decks from MDPR" for shape in slide.shapes)
        )
        overview_text_parts = [
            shape.text.lower()
            for shape in overview.shapes
            if getattr(shape, "has_text_frame", False)
        ]
        for shape in overview.shapes:
            if getattr(shape, "has_table", False):
                overview_text_parts.extend(cell.text.lower() for row in shape.table.rows for cell in row.cells)
        overview_text = "\n".join(overview_text_parts)
        overview_families = [
            match.removesuffix("-en")
            for match in re.findall(r"([a-z0-9-]+)/deck\.md", overview_text)
        ]

        self.assertEqual(len(selected_families), 6)
        self.assertEqual(len(selected_families), len(set(selected_families)))
        self.assertEqual(len(overview_lines), 3)
        self.assertIn("five-methods/deck.md — Five-Item Layout Example", overview_lines[-1])
        self.assertIn("theme-preview-en/deck.md", overview_lines[-1])
        self.assertEqual(overview_families, selected_families)
        for unselected in ("readme-final", "readme-teaser", "language-preview"):
            self.assertNotIn(unselected, overview_text)

    def test_source_family_samples_use_shortest_unique_suffix(self) -> None:
        module = load_script("create_mdpr_vs_skill_unique_source_labels", "scripts/create_mdpr_vs_skill_decks.py")
        self.assertEqual(
            module.shortest_unique_source_labels([
                "examples/basic/deck.md",
                "examples/comparison/deck.md",
                "examples/pipeline/deck.md",
            ]),
            ["basic/deck.md", "comparison/deck.md", "pipeline/deck.md"],
        )
        self.assertEqual(
            module.shortest_unique_source_labels([
                "docs/00-product-definition.md",
                "README.md",
                "docs/adr/0001-presentation-ir-schema-contract.md",
            ]),
            ["00-product-definition.md", "README.md", "0001-presentation-ir-schema-contract.md"],
        )
        self.assertEqual(
            module.shortest_unique_source_labels(["a/shared/file.md", "b/shared/file.md"]),
            ["a/shared/file.md", "b/shared/file.md"],
        )
        with self.assertRaises(ValueError):
            module.shortest_unique_source_labels(["C:/private/checkout/deck.md"])
        with self.assertRaises(ValueError):
            module.shortest_unique_source_labels(["../outside/deck.md"])
        with self.assertRaises(ValueError):
            module.shortest_unique_source_labels(["examples/basic/deck.md", "examples/basic/deck.md"])

        deck = Presentation(ROOT / "artifacts" / "mdpr-vs-skill" / "mdpr-skill-result.pptx")
        labels = [
            shape.text
            for shape in deck.slides[2].shapes
            if shape.name.startswith("group_1_item_")
        ]
        self.assertEqual(labels, ["basic/deck.md", "comparison/deck.md", "pipeline/deck.md"])

    def test_generated_mdpr_deck_has_no_exact_title_body_echo(self) -> None:
        module = load_script("create_mdpr_vs_skill_title_echo", "scripts/create_mdpr_vs_skill_decks.py")
        deck = Presentation(ROOT / "artifacts" / "mdpr-vs-skill" / "mdpr-baseline-result.pptx")
        echoes = []
        for slide_number, slide in enumerate(deck.slides, 1):
            texts = [
                shape.text.strip() for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            ]
            if not texts or texts[0] == "Agenda":
                continue
            title_key = module.heading_identity(texts[0], strip_continuation=True)
            duplicates = [text for text in texts[1:] if module.heading_identity(text) == title_key]
            if duplicates:
                echoes.append({"slide": slide_number, "title": texts[0], "duplicates": duplicates})

        self.assertEqual(echoes, [])

    def test_generated_agenda_ordinals_are_global_unique_and_contiguous(self) -> None:
        deck = Presentation(ROOT / "artifacts" / "mdpr-vs-skill" / "mdpr-baseline-result.pptx")
        agenda_items = []
        for slide in deck.slides:
            texts = [
                shape.text.strip() for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            ]
            if not texts or not texts[0].startswith("Agenda"):
                continue
            agenda_items.extend(text for text in texts[1:] if text[:2].isdigit() and text[2:4] == "  ")

        self.assertEqual([int(item[:2]) for item in agenda_items], list(range(1, 17)))
        self.assertEqual([item[4:] for item in agenda_items], [
            "Difference at a glance",
            "Pipeline boundary",
            "Architecture",
            "Page Splitting Rules",
            "Layout Selection Rules",
            "Rendering Rules",
            "Validation and Overflow Policy",
            "Example decks from MDPR",
            "Example: examples/basic/deck.md",
            "Example: examples/comparison/deck.md",
            "Example: examples/pipeline/deck.md",
            "Example: examples/diagram-arrangements/deck.md",
            "Example: examples/five-methods/deck.md",
            "Example: examples/theme-preview-en/deck.md",
            "Current skill output expectations",
            "End state",
        ])

    def test_exported_pngs_are_counted_once_on_case_insensitive_filesystems(self) -> None:
        module = load_script("create_mdpr_vs_skill_png_count", "scripts/create_mdpr_vs_skill_decks.py")
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir)
            (output / "Slide1.PNG").write_bytes(b"png")

            self.assertEqual(module.exported_png_paths(output), [output / "Slide1.PNG"])

    def test_powerpoint_export_uses_one_process_per_slide(self) -> None:
        module = load_script("create_mdpr_vs_skill_slide_export", "scripts/create_mdpr_vs_skill_decks.py")

        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            deck = Presentation()
            deck.slides.add_slide(deck.slide_layouts[6])
            deck.slides.add_slide(deck.slide_layouts[6])
            deck.save(source)
            output = root / "export"
            calls: list[tuple[list[str], dict]] = []

            def fake_run(command, **kwargs):
                calls.append((command, kwargs))
                path = Path(command[command.index("-OutputPath") + 1])
                write_palette_png(path)

            with patch.object(module.subprocess, "run", side_effect=fake_run):
                paths = module.export_with_powerpoint(source, output)

            self.assertEqual(paths, [output / "slide-001.png", output / "slide-002.png"])
            self.assertTrue(all(image_mode(path) == "RGB" for path in paths))
            self.assertEqual(len(calls), 2)
            commands = [call[0] for call in calls]
            self.assertEqual([call[call.index("-SlideIndex") + 1] for call in commands], ["1", "2"])
            self.assertTrue(all(kwargs["encoding"] == "utf-8" for _, kwargs in calls))
            self.assertTrue(all(kwargs["errors"] == "replace" for _, kwargs in calls))
            helper = (ROOT / "scripts" / "export_pptx_slide_isolated.ps1").read_text(encoding="utf-8")
            self.assertIn("GetWindowThreadProcessId", helper)
            self.assertIn("Wait-Process", helper)
            self.assertIn("Start-Sleep", helper)
            self.assertIn("warmup.png", helper)
            self.assertIn("Remove-Item -LiteralPath $warmupPath", helper)

    def test_powerpoint_export_retries_two_transient_slide_failures(self) -> None:
        module = load_script("create_mdpr_vs_skill_slide_retry", "scripts/create_mdpr_vs_skill_decks.py")
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            deck = Presentation()
            deck.slides.add_slide(deck.slide_layouts[6])
            deck.save(source)
            output = root / "export"
            attempts = 0

            def flaky_run(command, **_kwargs):
                nonlocal attempts
                attempts += 1
                if attempts < 3:
                    raise module.subprocess.CalledProcessError(1, command, stderr="transient")
                path = Path(command[command.index("-OutputPath") + 1])
                write_palette_png(path)

            with patch.object(module.subprocess, "run", side_effect=flaky_run):
                paths = module.export_with_powerpoint(source, output)

            self.assertEqual(paths, [output / "slide-001.png"])
            self.assertEqual(attempts, 3)

    def test_comparison_gate_rejects_failed_or_incomplete_current_exports(self) -> None:
        module = load_script("create_mdpr_vs_skill_report_gate", "scripts/create_mdpr_vs_skill_decks.py")
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            complete = complete_gate_report(root)

            self.assertTrue(module.comparison_report_ok(complete, actual_run_exists=True, artifact_root=root))
            complete["powerPointExport"]["ok"] = False
            self.assertFalse(module.comparison_report_ok(complete, actual_run_exists=True, artifact_root=root))
            complete["powerPointExport"]["ok"] = True
            complete["baselineRenderPreview"] = complete["baselineRenderPreview"][:3]
            self.assertFalse(module.comparison_report_ok(complete, actual_run_exists=True, artifact_root=root))

    def test_comparison_gate_rejects_sub_floor_typography_in_either_deck(self) -> None:
        module = load_script("create_mdpr_vs_skill_font_gate", "scripts/create_mdpr_vs_skill_decks.py")
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            complete = complete_gate_report(root)

            complete["mdprBaselineValidation"]["minFontSizePt"] = 15.9
            self.assertFalse(module.comparison_report_ok(complete, actual_run_exists=True, artifact_root=root))
            complete["mdprBaselineValidation"]["minFontSizePt"] = 16
            complete["skillValidation"]["minFontSizePt"] = 15.9
            self.assertFalse(module.comparison_report_ok(complete, actual_run_exists=True, artifact_root=root))

    def test_comparison_gate_rejects_named_card_content_beyond_its_container(self) -> None:
        module = load_script("create_mdpr_vs_skill_named_container_gate", "scripts/create_mdpr_vs_skill_decks.py")
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            complete = complete_gate_report(root)
            complete["skillValidation"]["namedContainerOverflowCount"] = 1

            self.assertFalse(module.comparison_report_ok(complete, actual_run_exists=True, artifact_root=root))


class ValidatePackCheckboxContractTests(unittest.TestCase):
    def test_intentional_github_form_checkboxes_are_not_unfinished_project_work(self) -> None:
        module = load_script("validate_pack", "scripts/validate_pack.py")
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            template = root / ".github" / "pull_request_template.md"
            template.parent.mkdir(parents=True)
            template.write_text("- [ ] I ran the tests.\n", encoding="utf-8")
            docs = root / "docs" / "done.md"
            docs.parent.mkdir(parents=True)
            docs.write_text("- [x] complete\n", encoding="utf-8")

            module.ROOT = root
            module.check_no_unchecked_boxes()

    def test_unchecked_boxes_in_governed_docs_still_fail(self) -> None:
        module = load_script("validate_pack_governed", "scripts/validate_pack.py")
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            docs = root / "docs" / "work.md"
            docs.parent.mkdir(parents=True)
            docs.write_text("- [ ] unfinished implementation\n", encoding="utf-8")

            module.ROOT = root
            with contextlib.redirect_stderr(io.StringIO()):
                with self.assertRaises(SystemExit):
                    module.check_no_unchecked_boxes()


if __name__ == "__main__":
    unittest.main()
