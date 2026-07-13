#!/usr/bin/env python3
from __future__ import annotations

import json
import os
import re
import shutil
import subprocess
import time
import unicodedata
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

try:
    from scripts.comparison_report_gate import MIN_FONT_SIZE_PT, comparison_report_ok
except ModuleNotFoundError:
    from comparison_report_gate import MIN_FONT_SIZE_PT, comparison_report_ok

ROOT = Path(__file__).resolve().parents[1]


def resolve_mdpr_root() -> Path:
    configured = os.environ.get("MDPR_SOURCE_DIR")
    candidates = [
        Path(configured).expanduser() if configured else None,
        ROOT.parent / "mdpresent-spec-scaffold",
        ROOT / ".cache" / "mdpr",
    ]
    for candidate in candidates:
        if candidate is not None and (candidate / "packages" / "cli").is_dir():
            return candidate.resolve()
    return (ROOT / ".cache" / "mdpr").resolve()


MDPR = resolve_mdpr_root()
OUT = ROOT / "artifacts" / "mdpr-vs-skill"
SOURCE_MD = OUT / "mdpr-source-corpus.md"
MANIFEST = OUT / "source-manifest.json"
BASELINE_PPTX = OUT / "mdpr-baseline-result.pptx"
SKILL_PPTX = OUT / "mdpr-skill-result.pptx"
SKILL_FROM_MDPR_RUN_PPTX = OUT / "mdpr-skill-from-actual-md-run.pptx"
REPORT = OUT / "mdpr-vs-skill-report.json"

SLIDE_W = 13.333
SLIDE_H = 7.5
SOURCE_FILES = [
    "README.md",
    "docs/00-product-definition.md",
    "docs/01-architecture.md",
    "docs/02-requirements.md",
    "docs/03-page-splitting.md",
    "docs/04-layout-rules.md",
    "docs/05-overrides-for-llm.md",
    "docs/06-cli-spec.md",
    "docs/07-rendering-rules.md",
    "docs/08-roadmap.md",
    "docs/09-codex-implementation-guide.md",
    "docs/10-template-and-master-policy.md",
    "docs/11-qa-overflow.md",
    "docs/references.md",
    "docs/adr/0001-presentation-ir-schema-contract.md",
    "examples/basic/deck.md",
    "examples/comparison/deck.md",
    "examples/pipeline/deck.md",
    "examples/diagram-arrangements/deck.md",
    "examples/five-methods/deck.md",
    "examples/theme-preview-en/deck.md",
]


@dataclass(frozen=True)
class Palette:
    bg: str = "F7F2EA"
    ink: str = "111827"
    muted: str = "5F6B7A"
    line: str = "D7CDBD"
    card: str = "FFFDF8"
    surface: str = "EAF7F3"
    surface2: str = "F1F0FA"
    accent: str = "0F766E"
    accent2: str = "7C3AED"
    contrast: str = "BE123C"
    amber: str = "A16207"
    dark: str = "20242C"


P = Palette()
MDPR_ACTOR_COLOR = P.accent
SKILL_ACTOR_COLOR = P.accent2


def evidence_path(path: Path) -> str:
    resolved = path.resolve()
    try:
        return str(resolved.relative_to(ROOT.resolve()))
    except ValueError:
        return str(resolved)


def rgb(hex_value: str) -> RGBColor:
    return RGBColor(*(int(hex_value[i:i + 2], 16) for i in (0, 2, 4)))


def clean_line(line: str) -> str:
    line = re.sub(r"!\[[^\]]*\]\([^)]*\)", "", line)
    line = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", line)
    line = line.replace("`", "")
    return line.strip()


def read_source(relative_path: str) -> str:
    return (MDPR / relative_path).read_text(encoding="utf-8", errors="ignore")


def extract_title(text: str, fallback: str) -> str:
    for line in text.splitlines():
        match = re.match(r"^#\s+(.+)$", line.strip())
        if match:
            return clean_line(match.group(1))[:80]
    return fallback


def extract_bullets(text: str, limit: int = 6) -> list[str]:
    bullets: list[str] = []
    for line in text.splitlines():
        stripped = clean_line(line)
        if re.match(r"^[-*]\s+\S", stripped):
            bullets.append(re.sub(r"^[-*]\s+", "", stripped)[:130])
        elif re.match(r"^\d+\.\s+\S", stripped):
            bullets.append(re.sub(r"^\d+\.\s+", "", stripped)[:130])
        if len(bullets) >= limit:
            break
    return bullets


def extract_code_block(text: str) -> tuple[str, list[str]] | None:
    lines = text.splitlines()
    for index, line in enumerate(lines):
        if not line.strip().startswith("```"):
            continue
        language = line.strip().strip("`") or "text"
        collected: list[str] = []
        for next_line in lines[index + 1:]:
            if next_line.strip().startswith("```"):
                break
            if next_line.strip():
                collected.append(next_line.rstrip())
            if len(collected) >= 8:
                break
        if collected:
            return language, collected
    return None


def extract_table(text: str) -> list[list[str]]:
    rows: list[list[str]] = []
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped.startswith("|") or not stripped.endswith("|"):
            continue
        cells = [clean_line(cell) for cell in stripped.strip("|").split("|")]
        if all(re.fullmatch(r":?-{3,}:?", cell.replace(" ", "")) for cell in cells):
            continue
        rows.append(cells[:4])
        if len(rows) >= 5:
            break
    return rows


def source_summary(relative_path: str) -> dict[str, Any]:
    text = read_source(relative_path)
    headings = [clean_line(match.group(1)) for match in re.finditer(r"^#{1,3}\s+(.+)$", text, flags=re.MULTILINE)]
    bullets = extract_bullets(text)
    table = extract_table(text)
    code = extract_code_block(text)
    return {
        "path": relative_path,
        "title": extract_title(text, Path(relative_path).stem),
        "headingCount": len(headings),
        "headings": headings[:8],
        "bullets": bullets,
        "hasTable": bool(table),
        "table": table,
        "hasCode": bool(code),
        "codeLanguage": code[0] if code else None,
        "code": code[1] if code else [],
        "charCount": len(text),
    }


def corpus_display_title(value: str) -> str:
    return re.sub(r"^\s*\d{1,3}[.)]\s+", "", value).strip()


def heading_identity(value: str, *, strip_continuation: bool = False) -> str:
    normalized = unicodedata.normalize("NFKC", value)
    if strip_continuation:
        normalized = re.sub(r"\s*\(continued(?:\s+\d+)?\)\s*$", "", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"^\s*\d{1,3}(?:(?:\s*[\W_]+\s*)|\s+)", "", normalized)
    normalized = "".join(character if character.isalnum() else " " for character in normalized.casefold())
    return " ".join(normalized.split())


def unique_corpus_headings(headings: list[str], section_title: str, *, limit: int) -> list[str]:
    seen = {heading_identity(section_title)}
    retained: list[str] = []
    for heading in headings:
        display_heading = corpus_display_title(heading)
        identity = heading_identity(display_heading)
        if not identity or identity in seen:
            continue
        seen.add(identity)
        retained.append(display_heading)
        if len(retained) >= limit:
            break
    return retained


def build_source_corpus(summaries: list[dict[str, Any]]) -> None:
    lines = [
        "# MDPR Corpus: Runtime vs Review Evidence",
        "",
        "This deck is generated from Markdown files inside the local MDPR checkout.",
        "",
        "## Difference at a glance",
        "",
        "| Area | MDPR runtime | mdpr-skill review companion |",
        "|---|---|---|",
        "| Role | Markdown to Presentation IR and rendered output | Optional semantic hints, review findings, and evidence |",
        "| Parser | Built-in parser or Pandoc parser mode | Does not parse Markdown; reads MDPR evidence |",
        "| Layout | Owns deterministic layout, typography, and theme rules | Does not set final coordinates, font sizes, or theme values |",
        "| Output | Editable PPTX, HTML, PDF, manifests, and previews | JSON hints, review reports, and comparison evidence |",
    ]
    lines.extend([
        "",
        "## Pipeline boundary",
        "",
        "Markdown => MDPR parser => BlockIR => Outline Tree => Split Planner => Presentation IR => Layout IR => Renderer",
        "",
        "MDPR manifest and previews => mdpr-skill hints or review findings => MDPR remains the only renderer",
        "",
    ])
    topic_paths = [
        "docs/01-architecture.md",
        "docs/03-page-splitting.md",
        "docs/04-layout-rules.md",
        "docs/07-rendering-rules.md",
        "docs/11-qa-overflow.md",
    ]
    for path in topic_paths:
        item = next(summary for summary in summaries if summary["path"] == path)
        section_title = corpus_display_title(item["title"])
        lines.extend([f"## {section_title}", ""])
        for heading in unique_corpus_headings(item["headings"], section_title, limit=5):
            lines.append(f"- {heading}")
        for bullet in item["bullets"][:4]:
            lines.append(f"- {bullet}")
        if item["table"]:
            lines.extend(["", "| Field | Value |", "|---|---|"])
            for row in item["table"][:4]:
                lines.append("| " + " | ".join(row[:2]).replace("\n", " ") + " |")
        if item["code"]:
            lines.extend(["", f"```{item['codeLanguage'] or 'text'}"])
            lines.extend(item["code"][:6])
            lines.append("```")
        lines.append("")
    lines.extend([
        "## Example decks from MDPR",
        "",
        "- basic/deck.md covers core flow and expected effects.",
        "- comparison/deck.md exercises before/after content.",
        "- pipeline/deck.md exercises diagram conversion.",
        "- diagram-arrangements/deck.md exercises multiple diagram structures.",
        "- theme-preview decks exercise preset variety.",
        "",
    ])
    for path in [item for item in summaries if item["path"].startswith("examples/")]:
        section_title = f"Example: {path['path']}"
        lines.extend([f"## {section_title}", ""])
        for heading in unique_corpus_headings(path["headings"], section_title, limit=6):
            lines.append(f"- {heading}")
        for bullet in path["bullets"][:4]:
            lines.append(f"- {bullet}")
        lines.append("")
    lines.extend([
        "## Current skill output expectations",
        "",
        "- Text-only slides may receive one quiet monotone-icon-aside slot.",
        "- Dense content should stay readable instead of gaining decorative icons.",
        "- Infographic families are selected by text length, relation, item count, and importance.",
        "- Coherence validation checks color role, alignment, object variety, font floor, and z-order.",
        "",
        "## End state",
        "",
        "> MDPR remains the content and rendering runtime. mdpr-skill adds optional semantic hints, review findings, and evidence without owning final visual decisions.",
    ])
    SOURCE_MD.write_text("\n".join(lines) + "\n", encoding="utf-8")


def build_mdpr_baseline() -> None:
    cli = MDPR / "packages" / "cli" / "dist" / "index.js"
    if not cli.is_file():
        raise FileNotFoundError(f"MDPR CLI build output is missing: {cli}")
    dist = OUT / "mdpr-baseline-build"
    if dist.exists():
        shutil.rmtree(dist)
    dist.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        ["node", str(cli), "build", str(SOURCE_MD), "--to", "pptx", "--out", str(dist), "--design", "clean"],
        cwd=MDPR,
        check=True,
    )
    generated = dist / "deck.pptx"
    if not generated.is_file():
        raise FileNotFoundError(f"MDPR did not write {generated}")
    shutil.copyfile(generated, BASELINE_PPTX)


def add_text(slide, name: str, x: float, y: float, w: float, h: float, text: str, size: int, color: str = P.ink, bold: bool = False, align: PP_ALIGN | None = None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(max(size, MIN_FONT_SIZE_PT))
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_shape(slide, name: str, kind: MSO_AUTO_SHAPE_TYPE, x: float, y: float, w: float, h: float, fill: str, line: str | None = None, text: str | None = None, font_size: int = 11, font_color: str = P.ink, bold: bool = False):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    if text is not None:
        tf = shape.text_frame
        tf.clear()
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(max(font_size, MIN_FONT_SIZE_PT))
        run.font.bold = bold
        run.font.color.rgb = rgb(font_color)
    return shape


def add_bg(slide) -> None:
    add_shape(slide, "z00_bg", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, P.bg, P.bg)


def draw_mono_icon(slide, name: str, x: float, y: float, size: float, kind: str, tone: str = P.ink) -> None:
    if kind == "pipeline":
        for i in range(3):
            add_shape(slide, f"{name}_node_{i}", MSO_AUTO_SHAPE_TYPE.OVAL, x + i * size * 0.38, y + size * 0.28, size * 0.18, size * 0.18, tone, tone)
        for i in range(2):
            add_shape(slide, f"{name}_line_{i}", MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + size * (0.17 + i * 0.38), y + size * 0.36, size * 0.22, size * 0.035, tone, tone)
    elif kind == "shield":
        add_shape(slide, f"{name}_body", MSO_AUTO_SHAPE_TYPE.PENTAGON, x, y, size * 0.55, size * 0.65, tone, tone)
        add_shape(slide, f"{name}_cut", MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + size * 0.22, y + size * 0.12, size * 0.1, size * 0.28, P.bg, P.bg)
    elif kind == "spark":
        add_shape(slide, f"{name}_v", MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + size * 0.24, y, size * 0.06, size * 0.58, tone, tone)
        add_shape(slide, f"{name}_h", MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y + size * 0.24, size * 0.58, size * 0.06, tone, tone)
        add_shape(slide, f"{name}_dot", MSO_AUTO_SHAPE_TYPE.OVAL, x + size * 0.44, y + size * 0.44, size * 0.1, size * 0.1, tone, tone)
    else:
        add_shape(slide, f"{name}_doc", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, size * 0.48, size * 0.62, tone, tone)
        for i in range(3):
            add_shape(slide, f"{name}_line_{i}", MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + size * 0.12, y + size * (0.16 + i * 0.13), size * 0.25, size * 0.035, P.bg, P.bg)


def add_title(slide, title: str) -> None:
    add_text(slide, "title", 0.65, 0.42, 11.9, 0.58, title, 28, P.ink, True)


def add_card(slide, name: str, x: float, y: float, w: float, h: float, title: str, body: list[str], accent: str = P.accent, icon: str | None = None) -> None:
    add_shape(slide, f"{name}_card", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, P.card, P.line)
    if icon:
        draw_mono_icon(slide, f"{name}_icon", x + 0.18, y + 0.18, 0.65, icon)
        tx = x + 0.88
        tw = w - 1.08
    else:
        tx = x + 0.22
        tw = w - 0.44
    add_text(slide, f"{name}_title", tx, y + 0.2, tw, 0.38, title, 18, accent, True)
    for i, line in enumerate(body[:3]):
        add_text(slide, f"{name}_body_{i}", tx, y + 0.7 + i * 0.42, tw, 0.34, line, 16, P.muted)


def add_difference_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "MDPR runtime and mdpr-skill boundary")
    headers = ["Area", "MDPR", "Current skill"]
    rows = [
        ["Core role", "Markdown -> editable presentation", "Optional hints and review evidence"],
        ["Parser", "Built-in or Pandoc mode", "No Markdown parsing"],
        ["Design", "Owns layout and theme binding", "Suggests intent; no final geometry"],
        ["Validation", "Owns deterministic pass/fail", "Mirrors MDPR findings with evidence"],
        ["Output", "Editable PPTX, HTML, and PDF", "JSON hints, reports, and evidence"],
    ]
    table = slide.shapes.add_table(len(rows) + 1, 3, Inches(0.68), Inches(1.62), Inches(11.9), Inches(4.1)).table
    widths = [1.7, 4.25, 5.95]
    for col, width in enumerate(widths):
        table.columns[col].width = Inches(width)
    for col, value in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = value
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(P.dark)
    for row_idx, row in enumerate(rows, 1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(P.surface if col_idx == 2 else P.card)
    for row in table.rows:
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(16)
                    run.font.color.rgb = rgb("FFFFFF" if cell in table.rows[0].cells else P.ink)
                    run.font.bold = cell in table.rows[0].cells


def add_actual_mdpr_run_slide(prs: Presentation, summaries: list[dict[str, Any]], mdpr_result: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Evidence starts with an actual MDPR run")
    add_card(
        slide,
        "actual_source",
        0.75,
        1.45,
        3.65,
        2.0,
        "1. Real Markdown input",
        [
            f"{len(summaries)} source files",
            f"{sum(item['headingCount'] for item in summaries)} headings",
            f"{sum(item['charCount'] for item in summaries):,} chars",
        ],
        P.accent,
        "doc",
    )
    add_card(
        slide,
        "actual_mdpr",
        4.85,
        1.45,
        3.65,
        2.0,
        "2. MDPR output",
        [
            f"{mdpr_result['slides']} slides",
            f"{mdpr_result['textFrames']} editable text frames",
            f"{mdpr_result['tables']} tables, {mdpr_result['charts']} charts",
        ],
        MDPR_ACTOR_COLOR,
        "pipeline",
    )
    add_card(
        slide,
        "actual_skill",
        8.95,
        1.45,
        3.65,
        2.0,
        "3. mdpr-skill review",
        [
            "reads MDPR evidence",
            "returns optional findings",
            "does not render final slides",
        ],
        SKILL_ACTOR_COLOR,
        "spark",
    )
    for i, x in enumerate([4.42, 8.52]):
        add_shape(slide, f"actual_arrow_{i}", MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, x, 2.3, 0.34, 0.28, P.dark, P.dark)
    add_shape(slide, "actual_boundary", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.85, 4.05, 11.55, 1.65, P.surface, P.line)
    add_text(slide, "actual_boundary_title", 1.12, 4.28, 3.1, 0.38, "Recorded boundary", 20, P.ink, True)
    add_text(
        slide,
        "actual_boundary_body",
        1.12,
        4.78,
        10.65,
        0.58,
        "MDPR creates the presentation. mdpr-skill may explain or review evidence, but it does not create a second styled runtime output.",
        16,
        P.muted,
    )


def add_source_coverage_slide(prs: Presentation, summaries: list[dict[str, Any]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "The same Markdown corpus grounds both sides")
    groups = [
        ("Docs", [s for s in summaries if s["path"].startswith("docs/") and "/adr/" not in s["path"]]),
        ("Examples", [s for s in summaries if s["path"].startswith("examples/")]),
        ("Root", [s for s in summaries if not s["path"].startswith(("docs/", "examples/"))]),
        ("ADR", [s for s in summaries if "/adr/" in s["path"]]),
    ]
    for i, (name, items) in enumerate(groups):
        x = 0.72 + i * 3.05
        add_shape(slide, f"group_{i}_card", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, 1.48, 2.65, 4.75, P.card, P.line)
        add_text(slide, f"group_{i}_num", x + 0.22, 1.78, 1.2, 0.55, str(len(items)), 30, P.accent2 if i == 1 else P.accent, True)
        add_text(slide, f"group_{i}_label", x + 0.24, 2.48, 1.8, 0.38, name, 20, P.ink, True)
        add_text(slide, f"group_{i}_chars", x + 0.24, 3.02, 2.1, 0.34, f"{sum(item['charCount'] for item in items):,} chars", 16, P.muted, True)
        for j, item in enumerate(items[:3]):
            add_text(slide, f"group_{i}_item_{j}", x + 0.24, 3.7 + j * 0.62, 2.08, 0.5, Path(item["path"]).name, 16, P.muted)


def add_pipeline_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "One runtime, one optional review companion")
    nodes = [
        ("Markdown", "source", P.card),
        ("Parser", "simple/Pandoc", P.card),
        ("Presentation IR", "semantic slides", P.surface),
        ("MDPR rules", "layout/theme", P.surface),
        ("Editable PPTX", "runtime output", P.card),
    ]
    for i, (title, sub, fill) in enumerate(nodes):
        x = 0.72 + i * 2.45
        add_shape(slide, f"node_{i}", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, 2.25, 1.82, 1.12, fill, P.line)
        add_text(slide, f"node_{i}_title", x + 0.12, 2.4, 1.58, 0.42, title, 16, P.ink, True, PP_ALIGN.CENTER)
        add_text(slide, f"node_{i}_sub", x + 0.12, 2.86, 1.58, 0.34, sub, 16, P.muted, False, PP_ALIGN.CENTER)
        if i < len(nodes) - 1:
            add_shape(slide, f"arrow_{i}", MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, x + 1.86, 2.62, 0.48, 0.32, P.accent if i < 2 else P.accent2, P.accent if i < 2 else P.accent2)
    add_card(slide, "mdpr", 0.82, 4.08, 5.35, 2.05, "MDPR owns", ["layout and typography", "editable rendering", "deterministic pass/fail"], MDPR_ACTOR_COLOR, "doc")
    add_card(slide, "skill", 6.72, 4.08, 5.35, 2.05, "mdpr-skill assists", ["semantic hints", "review explanations", "evidence routing"], SKILL_ACTOR_COLOR, "spark")


def add_docs_map_slide(prs: Presentation, summaries: list[dict[str, Any]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Source families exercise different runtime paths")
    rows = [
        ("Architecture", "packages and renderer", "pipeline evidence"),
        ("Page splitting", "headings and density", "split-rule evidence"),
        ("Layout rules", "regions and overflow", "grid evidence"),
        ("Rendering rules", "PPTX, PDF, and HTML", "format evidence"),
        ("Overrides", "human and LLM boundary", "annotation evidence"),
        ("Examples", "actual deck Markdown", "intent evidence"),
    ]
    for i, (title, body, output) in enumerate(rows):
        x = 0.75 + (i % 3) * 4.0
        y = 1.48 + (i // 3) * 2.25
        add_card(slide, f"docmap_{i}", x, y, 3.45, 1.9, title, [body, output], P.accent if i % 2 == 0 else P.accent2, ["doc", "pipeline", "shield"][i % 3])


def add_visual_rules_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "mdpr-skill narrows review, not rendering")
    add_card(slide, "icon", 0.78, 1.48, 3.65, 1.92, "Semantic hints", ["grouping and intent", "source-bound keywords", "content split candidates"], P.ink, "spark")
    add_card(slide, "info", 4.86, 1.48, 3.65, 1.92, "Review evidence", ["rendered preview paths", "MDPR finding IDs", "manifest status"], P.accent, "pipeline")
    add_card(slide, "qa", 8.94, 1.48, 3.65, 1.92, "Hard boundary", ["no coordinates", "no exact typography", "no pass/fail override"], P.contrast, "shield")
    add_shape(slide, "line", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 1.05, 4.03, 10.95, 0.04, P.line, P.line)
    steps = [("1", "Detect text-only"), ("2", "Reserve aside slot"), ("3", "Pick monotone source"), ("4", "Validate alignment")]
    for i, (num, label) in enumerate(steps):
        x = 1.1 + i * 2.9
        add_shape(slide, f"step_{i}_num", MSO_AUTO_SHAPE_TYPE.OVAL, x, 3.76, 0.46, 0.46, P.dark, P.dark, num, 16, "FFFFFF", True)
        add_text(slide, f"step_{i}_label", x - 0.64, 4.4, 1.75, 0.56, label, 16, P.ink, True, PP_ALIGN.CENTER)


def add_examples_slide(prs: Presentation, summaries: list[dict[str, Any]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Examples broaden content without changing ownership")
    examples = [s for s in summaries if s["path"].startswith("examples/")][:6]
    for i, item in enumerate(examples):
        x = 0.72 + (i % 3) * 4.0
        y = 1.48 + (i // 3) * 2.25
        family_names = {
            "diagram-arrangements": "Diagram",
            "five-methods": "Five Methods",
            "theme-preview-en": "Theme Preview",
        }
        directory = Path(item["path"]).parent.name
        family = family_names.get(directory, directory.replace("-", " ").title())
        body = [item["title"][:45], f"{item['headingCount']} headings · {item['charCount']:,} chars"]
        add_card(slide, f"example_{i}", x, y, 3.45, 1.9, family, body, P.accent2 if i % 2 else P.accent, "doc")


def add_chart_slide(prs: Presentation, summaries: list[dict[str, Any]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Corpus shape is evidence, not a second renderer")
    docs = [s for s in summaries if s["path"].startswith("docs/")]
    examples = [s for s in summaries if s["path"].startswith("examples/")]
    root = [s for s in summaries if not s["path"].startswith(("docs/", "examples/"))]
    data = CategoryChartData()
    data.categories = ["Docs", "Examples", "Root"]
    data.add_series("Headings", (sum(s["headingCount"] for s in docs), sum(s["headingCount"] for s in examples), sum(s["headingCount"] for s in root)))
    chart_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.92), Inches(1.75), Inches(5.35), Inches(3.4), data)
    chart_frame.name = "native_heading_chart"
    chart = chart_frame.chart
    chart.has_title = False
    chart.has_legend = False
    chart.value_axis.tick_labels.font.size = Pt(16)
    chart.category_axis.tick_labels.font.size = Pt(16)
    table = slide.shapes.add_table(4, 4, Inches(6.85), Inches(1.82), Inches(5.15), Inches(3.05)).table
    rows = [["Group", "Files", "Headings", "Chars"], ["Docs", str(len(docs)), str(sum(s["headingCount"] for s in docs)), f"{sum(s['charCount'] for s in docs):,}"], ["Examples", str(len(examples)), str(sum(s["headingCount"] for s in examples)), f"{sum(s['charCount'] for s in examples):,}"], ["Root", str(len(root)), str(sum(s["headingCount"] for s in root)), f"{sum(s['charCount'] for s in root):,}"]]
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(P.surface if r == 0 else P.card)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(16)
                    run.font.bold = r == 0
                    run.font.color.rgb = rgb(P.ink)


def add_text_icon_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Suggestions stop before geometry")
    add_text(slide, "suggestion_heading", 0.88, 1.55, 5.25, 0.45, "mdpr-skill may suggest", 20, SKILL_ACTOR_COLOR, True)
    suggestions = [
        ("PLACEMENT", "Reserve an aside or corner region."),
        ("SOURCE", "Use one licensed monochrome icon."),
        ("CHECK", "Compare visual center with the text midpoint."),
    ]
    for i, (label, body) in enumerate(suggestions):
        y = 2.25 + i * 1.18
        add_shape(slide, f"suggestion_{i}_mark", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.9, y + 0.05, 0.18, 0.58, SKILL_ACTOR_COLOR, SKILL_ACTOR_COLOR)
        add_text(slide, f"suggestion_{i}_label", 1.28, y, 1.25, 0.3, label, 16, P.ink, True)
        add_text(slide, f"suggestion_{i}_body", 2.58, y - 0.02, 3.55, 0.68, body, 16, P.muted)

    add_shape(slide, "runtime_field", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 6.72, 1.42, 5.9, 4.95, P.surface, P.surface)
    add_text(slide, "runtime_heading", 7.18, 1.55, 4.85, 0.45, "MDPR still decides", 20, MDPR_ACTOR_COLOR, True)
    decisions = [
        ("01", "Exact geometry", "regions and coordinates"),
        ("02", "Typography", "font family, size, and wrapping"),
        ("03", "Pass or fail", "deterministic validation result"),
    ]
    for i, (num, title, body) in enumerate(decisions):
        y = 2.22 + i * 1.2
        add_text(slide, f"decision_{i}_num", 7.18, y, 0.58, 0.42, num, 18, MDPR_ACTOR_COLOR, True)
        add_text(slide, f"decision_{i}_title", 7.92, y, 3.9, 0.34, title, 18, P.ink, True)
        add_text(slide, f"decision_{i}_body", 7.92, y + 0.42, 3.9, 0.34, body, 16, P.muted)


def build_skill_deck(summaries: list[dict[str, Any]], mdpr_result: dict[str, Any]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slides = [
        lambda: add_actual_mdpr_run_slide(prs, summaries, mdpr_result),
        lambda: add_difference_slide(prs),
        lambda: add_source_coverage_slide(prs, summaries),
        lambda: add_pipeline_slide(prs),
        lambda: add_docs_map_slide(prs, summaries),
        lambda: add_visual_rules_slide(prs),
        lambda: add_examples_slide(prs, summaries),
        lambda: add_chart_slide(prs, summaries),
        lambda: add_text_icon_slide(prs),
    ]
    for add in slides:
        add()
    prs.save(SKILL_PPTX)
    shutil.copyfile(SKILL_PPTX, SKILL_FROM_MDPR_RUN_PPTX)


def exported_png_paths(output_dir: Path) -> list[Path]:
    return sorted(path for path in output_dir.iterdir() if path.is_file() and path.suffix.lower() == ".png")


def wait_for_stable_export(path: Path, *, timeout_seconds: float = 10.0, settle_seconds: float = 0.2) -> None:
    deadline = time.monotonic() + timeout_seconds
    previous_size = -1
    stable_checks = 0
    while time.monotonic() < deadline:
        if path.is_file():
            size = path.stat().st_size
            if size > 0 and size == previous_size:
                stable_checks += 1
                if stable_checks >= 2:
                    if settle_seconds:
                        time.sleep(settle_seconds)
                    return
            else:
                stable_checks = 0
            previous_size = size
        time.sleep(0.05)
    raise TimeoutError(f"PowerPoint did not finish exporting {path}")


def export_with_powerpoint(pptx_path: Path, output_dir: Path, width: int = 1600, height: int = 900) -> list[Path]:
    if output_dir.exists():
        shutil.rmtree(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    slide_count = len(Presentation(pptx_path).slides)
    helper = ROOT / "scripts" / "export_pptx_slide_isolated.ps1"
    exported: list[Path] = []
    for index in range(1, slide_count + 1):
        output = output_dir / f"slide-{index:03d}.png"
        command = [
            "powershell",
            "-NoProfile",
            "-NonInteractive",
            "-ExecutionPolicy",
            "Bypass",
            "-File",
            str(helper),
            "-PptxPath",
            str(pptx_path.resolve()),
            "-OutputPath",
            str(output.resolve()),
            "-SlideIndex",
            str(index),
            "-Width",
            str(width),
            "-Height",
            str(height),
        ]
        for attempt in range(2):
            try:
                subprocess.run(
                    command,
                    cwd=ROOT,
                    check=True,
                    capture_output=True,
                    text=True,
                    encoding="utf-8",
                    errors="replace",
                )
                break
            except subprocess.CalledProcessError:
                output.unlink(missing_ok=True)
                if attempt == 1:
                    raise
        wait_for_stable_export(output, settle_seconds=0)
        exported.append(output)
    return exported


def summarize_export(paths: list[Path]) -> dict[str, Any]:
    validation = validate_pngs(paths)
    invalid = [item["file"] for item in validation if not item["hasContent"]]
    return {"count": len(paths), "invalidSlideCount": len(invalid), "invalidSlides": invalid}


def validate_pptx(path: Path) -> dict[str, Any]:
    prs = Presentation(path)
    counts = {"slides": len(prs.slides), "shapes": 0, "textFrames": 0, "pictures": 0, "tables": 0, "charts": 0}
    min_font = 999.0
    named_container_overflow: list[str] = []
    for slide_index, slide in enumerate(prs.slides, 1):
        named_shapes = {shape.name: shape for shape in slide.shapes}
        for name, container in named_shapes.items():
            if not name.endswith("_card"):
                continue
            prefix = name[:-5]
            children = [shape for child_name, shape in named_shapes.items() if child_name.startswith(f"{prefix}_body_")]
            if any(child.top + child.height > container.top + container.height for child in children):
                named_container_overflow.append(f"slide-{slide_index}:{name}")
        for shape in slide.shapes:
            counts["shapes"] += 1
            if getattr(shape, "has_text_frame", False):
                counts["textFrames"] += 1
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is not None and run.text.strip():
                            min_font = min(min_font, float(run.font.size.pt))
            if getattr(shape, "has_table", False):
                counts["tables"] += 1
            if getattr(shape, "has_chart", False):
                counts["charts"] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                counts["pictures"] += 1
    return {
        **counts,
        "minFontSizePt": None if min_font == 999.0 else min_font,
        "namedContainerOverflowCount": len(named_container_overflow),
        "namedContainerOverflow": named_container_overflow,
    }


def validate_pngs(paths: list[Path]) -> list[dict[str, Any]]:
    results = []
    for path in paths:
        image = Image.open(path).convert("RGB")
        colors = image.getcolors(maxcolors=4_000_000) or []
        non_white = sum(count for count, color in colors if color != (255, 255, 255))
        results.append({
            "file": str(path.relative_to(ROOT)),
            "size": image.size,
            "uniqueColors": len(colors),
            "nonWhitePixels": non_white,
            "hasContent": len(colors) > 20 and non_white > 50_000,
        })
    return results


def normalized_error_message(error: Exception) -> str:
    message = str(error)
    return message.replace("예외가 발생했습니다.", "PowerPoint COM exception")


def clear_preview_files(output_dir: Path) -> None:
    for pattern in ("mdpr_baseline_preview_*.png", "skill_preview_*.png"):
        for path in output_dir.glob(pattern):
            path.unlink()


def git_commit(repo: Path) -> str | None:
    result = subprocess.run(
        ["git", "rev-parse", "HEAD"],
        cwd=repo,
        check=False,
        capture_output=True,
        text=True,
    )
    return result.stdout.strip() or None if result.returncode == 0 else None


def main() -> None:
    if not MDPR.is_dir():
        raise FileNotFoundError("MDPR checkout is missing. Run npm run install:mdpr first.")
    OUT.mkdir(parents=True, exist_ok=True)
    summaries = [source_summary(path) for path in SOURCE_FILES if (MDPR / path).is_file()]
    MANIFEST.write_text(json.dumps({"sourceRoot": evidence_path(MDPR), "files": summaries}, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    build_source_corpus(summaries)
    build_mdpr_baseline()
    mdpr_baseline_validation = validate_pptx(BASELINE_PPTX)
    build_skill_deck(summaries, mdpr_baseline_validation)
    clear_preview_files(OUT)
    export_errors: list[dict[str, str]] = []
    try:
        baseline_exports = export_with_powerpoint(BASELINE_PPTX, OUT / "mdpr-baseline-export")
    except Exception as error:
        baseline_exports = []
        export_errors.append({"deck": str(BASELINE_PPTX.relative_to(ROOT)), "error": normalized_error_message(error)})
    try:
        skill_exports = export_with_powerpoint(SKILL_PPTX, OUT / "skill-export")
    except Exception as error:
        skill_exports = []
        export_errors.append({"deck": str(SKILL_PPTX.relative_to(ROOT)), "error": normalized_error_message(error)})
    for index, exported in enumerate(baseline_exports[:4], 1):
        shutil.copyfile(exported, OUT / f"mdpr_baseline_preview_{index}.png")
    for index, exported in enumerate(skill_exports[:4], 1):
        shutil.copyfile(exported, OUT / f"skill_preview_{index}.png")
    baseline_export_validation = summarize_export(baseline_exports)
    skill_export_validation = summarize_export(skill_exports)
    for temp_dir in [OUT / "mdpr-baseline-build", OUT / "mdpr-baseline-export", OUT / "skill-export"]:
        if temp_dir.exists():
            shutil.rmtree(temp_dir)
    report = {
        "sourceMarkdown": str(SOURCE_MD.relative_to(ROOT)),
        "sourceManifest": str(MANIFEST.relative_to(ROOT)),
        "mdprBaselinePptx": str(BASELINE_PPTX.relative_to(ROOT)),
        "skillPptx": str(SKILL_PPTX.relative_to(ROOT)),
        "skillFromActualMdprRunPptx": str(SKILL_FROM_MDPR_RUN_PPTX.relative_to(ROOT)),
        "sourceFileCount": len(summaries),
        "sourceHeadingCount": sum(item["headingCount"] for item in summaries),
        "sourceCharCount": sum(item["charCount"] for item in summaries),
        "actualMarkdownRun": {
            "inputMarkdown": str(SOURCE_MD.relative_to(ROOT)),
            "mdprSourceRoot": evidence_path(MDPR),
            "mdprCommit": git_commit(MDPR),
            "mdprCommand": f"node {evidence_path(MDPR / 'packages/cli/dist/index.js')} build {SOURCE_MD.relative_to(ROOT)} --to pptx --design clean",
            "mdprResultPptx": str(BASELINE_PPTX.relative_to(ROOT)),
            "mdprResultValidation": mdpr_baseline_validation,
            "skillResultPptx": str(SKILL_FROM_MDPR_RUN_PPTX.relative_to(ROOT)),
            "skillConsumes": [
                "the generated source Markdown corpus",
                "the source manifest extracted from the MDPR checkout",
                "the concrete MDPR PPTX run metrics",
            ],
        },
        "mdprBaselineValidation": mdpr_baseline_validation,
        "skillValidation": validate_pptx(SKILL_PPTX),
        "powerPointExport": {
            "ok": not export_errors,
            "errors": export_errors,
            "baselineSlideCount": len(baseline_exports),
            "skillSlideCount": len(skill_exports),
            "baselineValidation": baseline_export_validation,
            "skillValidation": skill_export_validation,
            "fallback": None,
            "evidencePolicy": "Each slide uses an isolated PowerPoint process, a stabilization delay, and a discarded warm-up frame; previews are cleared before each run and failed exports cannot reuse stale PNG evidence.",
        },
        "baselineRenderPreview": validate_pngs([path for path in [OUT / f"mdpr_baseline_preview_{index}.png" for index in range(1, 5)] if path.is_file()]),
        "skillRenderPreview": validate_pngs([path for path in [OUT / f"skill_preview_{index}.png" for index in range(1, 5)] if path.is_file()]),
    }
    report["ok"] = comparison_report_ok(report, actual_run_exists=SKILL_FROM_MDPR_RUN_PPTX.is_file())
    REPORT.write_text(json.dumps(report, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    if not report["ok"]:
        raise SystemExit(json.dumps(report, indent=2, ensure_ascii=False))
    print(json.dumps(report, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
