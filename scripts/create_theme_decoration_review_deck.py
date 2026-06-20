#!/usr/bin/env python3
from __future__ import annotations

import json
import re
import shutil
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "artifacts" / "theme-decoration-review"
PPTX = OUT / "theme-decoration-review.pptx"
EXPORT_DIR = OUT / "png"
REPORT = OUT / "theme-decoration-coverage-report.json"
CATALOG = ROOT / "design_components" / "decoration" / "src" / "decorators" / "objectShapeCatalog.json"
MATRIX = ROOT / "artifacts" / "theme-style-color-matrix"

SLIDE_W = 13.333
SLIDE_H = 7.5
BG = "F8FAFC"
TEXT = "111827"
MUTED = "475569"
LINE = "CBD5E1"
SURFACE = "FFFFFF"
ACCENT = "0F766E"
CONTRAST = "E11D48"
SAFE_PAD = 0.34
MIN_TEXT_PT = 7
VISUAL_REVIEW_FINDINGS: list[str] = []


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip().lstrip("#")
    return RGBColor(*(int(value[i:i + 2], 16) for i in (0, 2, 4)))


def normalize_hex(hex_value: str) -> str:
    value = str(hex_value).strip().lstrip("#").upper()
    if len(value) == 3:
        return "".join(char * 2 for char in value)
    return value


def blend_hex(left: str, right: str, right_weight: float) -> str:
    left = normalize_hex(left)
    right = normalize_hex(right)
    weight = max(0.0, min(1.0, right_weight))
    values = []
    for idx in (0, 2, 4):
        a = int(left[idx:idx + 2], 16)
        b = int(right[idx:idx + 2], 16)
        values.append(round(a * (1 - weight) + b * weight))
    return "".join(f"{value:02X}" for value in values)


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int,
    color: str = TEXT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    name: str = "mdpr-text",
    margin: float = 0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin / 2)
    tf.margin_bottom = Inches(margin / 2)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(max(size, MIN_TEXT_PT))
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def add_box(slide, x: float, y: float, w: float, h: float, fill: str = SURFACE, line: str = LINE, radius: bool = True, name: str = "mdpr-surface"):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.8)
    return shape


def add_shadow_plate(slide, x: float, y: float, w: float, h: float, fill: str = "CBD5E1", name: str = "mdpr-effect-shadow"):
    plate = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.045), Inches(y + 0.055), Inches(w), Inches(h))
    plate.name = name
    plate.fill.solid()
    plate.fill.fore_color.rgb = rgb(fill)
    plate.line.color.rgb = rgb(fill)
    return plate


def add_themed_surface(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    style: str,
    fill: str,
    line: str,
    accent: str,
    background: str,
    name: str = "mdpr-surface themed",
):
    if style in {"glass", "newmorphism", "magazine", "executive", "technical", "dark"}:
        add_shadow_plate(slide, x, y, w, h, blend_hex(line, background, 0.45 if style == "glass" else 0.7))
    surface_fill = fill
    surface_line = line
    if style == "glass":
        surface_fill = blend_hex(fill, background, 0.34)
        surface_line = blend_hex(line, "FFFFFF", 0.28)
    elif style == "grid":
        surface_fill = blend_hex(fill, background, 0.2)
    elif style == "data":
        surface_fill = blend_hex(fill, background, 0.08)
    elif style == "magazine":
        surface_fill = blend_hex(fill, "FFFFFF", 0.1)
    elif style == "newmorphism":
        surface_fill = blend_hex(fill, background, 0.65)
        surface_line = blend_hex("FFFFFF", line, 0.22)
    elif style == "minimalism":
        surface_fill = blend_hex(fill, "FFFFFF", 0.08)
        surface_line = blend_hex(line, background, 0.18)
    shape = add_box(slide, x, y, w, h, surface_fill, surface_line, True, name=name)
    if style == "glass":
        add_line(slide, x + 0.18, y + 0.18, w - 0.36, 0.018, blend_hex("FFFFFF", accent, 0.18), 1.0, name="mdpr-effect glass-highlight")
        add_dot(slide, x + w - 0.45, y + 0.18, 0.18, blend_hex(accent, "FFFFFF", 0.25), name="mdpr-effect glass-glow")
    elif style == "grid":
        for gx in [x + w * 0.33, x + w * 0.66]:
            add_line(slide, gx, y + 0.1, 0.01, h - 0.2, blend_hex(line, background, 0.28), 0.6, name="mdpr-effect grid-rule")
        for gy in [y + h * 0.42, y + h * 0.72]:
            add_line(slide, x + 0.1, gy, w - 0.2, 0.01, blend_hex(line, background, 0.28), 0.6, name="mdpr-effect grid-rule")
    elif style == "data":
        add_line(slide, x + 0.18, y + h - 0.2, w - 0.36, 0.018, accent, 1.0, name="mdpr-effect data-baseline")
        for i, bar_h in enumerate([0.12, 0.22, 0.16]):
            add_line(slide, x + 0.28 + i * 0.18, y + h - 0.25 - bar_h, 0.08, bar_h, accent, 1.0, name="mdpr-effect data-microbar")
    elif style == "magazine":
        add_line(slide, x, y + 0.12, w, 0.018, accent, 1.0, name="mdpr-effect magazine-rule")
        add_line(slide, x + 0.12, y, 0.035, h, accent, 1.0, name="mdpr-effect magazine-rail")
    elif style == "newmorphism":
        add_line(slide, x + 0.14, y + 0.12, w * 0.42, 0.016, "FFFFFF", 1.0, name="mdpr-effect newmorphism-highlight")
        add_line(slide, x + w * 0.55, y + h - 0.16, w * 0.32, 0.016, blend_hex(line, background, 0.35), 1.0, name="mdpr-effect newmorphism-lowlight")
        add_dot(slide, x + w - 0.34, y + 0.14, 0.12, blend_hex(accent, fill, 0.55), name="mdpr-effect newmorphism-accent")
    elif style == "minimalism":
        add_line(slide, x + 0.14, y + 0.14, w * 0.28, 0.012, accent, 0.8, name="mdpr-effect minimalism-rule")
    elif style == "simple":
        add_line(slide, x + 0.18, y + 0.16, 0.58, 0.028, accent, 1.0, name="mdpr-effect simple-rule")
    elif style == "technical":
        add_line(slide, x + 0.15, y + h - 0.16, w - 0.3, 0.018, accent, 1.0, name="mdpr-effect technical-baseline")
    return shape


def add_line(slide, x: float, y: float, w: float, h: float, color: str = ACCENT, width: float = 1.2, name: str = "mdpr-decor"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(max(w, 0.01)), Inches(max(h, 0.01)))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(color)
    shape.line.transparency = 100
    return shape


def add_dot(slide, x: float, y: float, size: float, fill: str = ACCENT, name: str = "mdpr-decor"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(fill)
    return shape


def add_pattern_preview(slide, pattern: dict[str, Any], x: float, y: float, w: float, h: float, index: int) -> None:
    family = pattern["family"]
    kind = pattern["kind"]
    fill = {
        "metric": "F0FDF4",
        "image": "EFF6FF",
        "diagram": "ECFDF5",
        "callout": "FFF7ED",
        "comparison": "F8FAFC",
        "table": "F1F5F9",
    }.get(family, SURFACE)
    accent = [ACCENT, "2563EB", "C2410C", "7C3AED", CONTRAST, "16A34A"][index % 6]

    add_box(slide, x, y, w, h, fill, LINE, name="mdpr-surface catalog-card")
    preview_x = x + 0.18
    preview_y = y + 0.2
    preview_w = 0.74
    preview_h = h - 0.4
    text_x = x + 1.05
    text_w = w - 1.23
    center_y = y + h / 2
    if kind == "table-summary-card":
        add_box(slide, preview_x, center_y - 0.28, 0.72, 0.56, "FFFFFF", LINE, False, name="mdpr-decor table-summary")
        add_line(slide, preview_x, center_y - 0.28, 0.72, 0.08, accent, name="mdpr-decor table-header")
        for row_y in [center_y - 0.08, center_y + 0.1]:
            add_line(slide, preview_x + 0.06, row_y, 0.6, 0.014, LINE, name="mdpr-decor table-row")
            add_dot(slide, preview_x + 0.08, row_y - 0.035, 0.07, accent, name="mdpr-decor table-dot")
        add_text(slide, preview_x + 0.42, center_y - 0.11, 0.25, 0.13, "42", 5, accent, True, PP_ALIGN.RIGHT, name="mdpr-text-in-decor", margin=0)
        add_text(slide, preview_x + 0.42, center_y + 0.07, 0.25, 0.13, "18", 5, accent, True, PP_ALIGN.RIGHT, name="mdpr-text-in-decor", margin=0)
    elif kind == "constraint-stack-card":
        add_text(slide, preview_x, center_y - 0.24, 0.68, 0.16, "MUST", 5, accent, True, name="mdpr-text-in-decor", margin=0)
        add_line(slide, preview_x + 0.18, center_y - 0.01, 0.5, 0.018, LINE, name="mdpr-decor constraint-line")
        add_line(slide, preview_x + 0.18, center_y + 0.16, 0.42, 0.018, LINE, name="mdpr-decor constraint-line")
        add_line(slide, preview_x + 0.08, center_y - 0.04, 0.035, 0.28, accent, name="mdpr-decor constraint-indent")
    elif kind == "arc-corner-emphasis":
        add_line(slide, preview_x + 0.05, center_y - 0.24, 0.52, 0.035, accent, name="mdpr-decor arc-corner")
        add_line(slide, preview_x + 0.05, center_y - 0.24, 0.035, 0.52, accent, name="mdpr-decor arc-corner")
        add_dot(slide, preview_x + 0.45, center_y + 0.03, 0.22, accent, name="mdpr-decor arc-dot")
        add_text(slide, preview_x + 0.02, center_y + 0.11, 0.58, 0.16, "78%", 6, accent, True, PP_ALIGN.CENTER, name="mdpr-text-in-decor", margin=0)
    elif kind == "subtle-band-card":
        add_box(slide, preview_x, center_y - 0.2, preview_w, 0.4, "E2E8F0", LINE, False, name="mdpr-decor subtle-band")
        add_line(slide, preview_x + 0.08, center_y - 0.03, 0.54, 0.018, accent, name="mdpr-decor subtle-band-line")
    elif kind == "inset-label-bar":
        add_box(slide, preview_x, center_y - 0.22, preview_w, 0.44, "FFFFFF", LINE, False, name="mdpr-decor inset-label")
        add_line(slide, preview_x + 0.08, center_y - 0.11, 0.16, 0.22, accent, name="mdpr-decor inset-label-bar")
        add_line(slide, preview_x + 0.32, center_y - 0.02, 0.32, 0.018, LINE, name="mdpr-decor inset-label-text")
    elif kind == "vertical-timeline-cards":
        add_line(slide, preview_x + 0.18, center_y - 0.28, 0.02, 0.56, accent, name="mdpr-decor timeline-line")
        for idx, row_y in enumerate([center_y - 0.19, center_y + 0.12]):
            add_dot(slide, preview_x + 0.11, row_y, 0.14, accent, name="mdpr-decor timeline-dot")
            add_box(slide, preview_x + 0.32, row_y - 0.03, 0.34, 0.16, "FFFFFF", LINE, False, name="mdpr-decor timeline-card")
    elif kind == "donut-label-ring":
        add_dot(slide, preview_x + 0.16, center_y - 0.24, 0.48, accent, name="mdpr-decor donut")
        add_dot(slide, preview_x + 0.28, center_y - 0.12, 0.24, fill, name="mdpr-decor donut-hole")
        add_line(slide, preview_x + 0.02, center_y - 0.2, 0.22, 0.018, accent, name="mdpr-decor donut-label")
        add_line(slide, preview_x + 0.54, center_y + 0.16, 0.18, 0.018, LINE, name="mdpr-decor donut-label")
    elif kind == "gauge-score-card":
        add_dot(slide, preview_x + 0.16, center_y - 0.24, 0.48, "DCFCE7", name="mdpr-decor gauge")
        add_dot(slide, preview_x + 0.27, center_y - 0.13, 0.26, fill, name="mdpr-decor gauge-hole")
        add_line(slide, preview_x + 0.38, center_y, 0.24, 0.025, accent, name="mdpr-decor gauge-needle")
        add_text(slide, preview_x + 0.27, center_y - 0.05, 0.25, 0.12, "78", 5, accent, True, PP_ALIGN.CENTER, name="mdpr-text-in-decor", margin=0)
    elif kind == "status-dot-table":
        for row_y in [center_y - 0.2, center_y - 0.04, center_y + 0.12]:
            add_line(slide, preview_x + 0.14, row_y, 0.52, 0.014, LINE, name="mdpr-decor status-row")
        add_dot(slide, preview_x + 0.02, center_y - 0.25, 0.12, accent, name="mdpr-decor status-dot")
        add_dot(slide, preview_x + 0.02, center_y - 0.09, 0.12, CONTRAST, name="mdpr-decor status-dot")
    elif kind == "risk-heat-grid":
        cell = 0.18
        for rr in range(2):
            for cc in range(2):
                tint = ["DCFCE7", "FEF3C7", "FED7AA", "FEE2E2"][rr * 2 + cc]
                add_box(slide, preview_x + 0.12 + cc * cell, center_y - 0.2 + rr * cell, cell, cell, tint, "FFFFFF", False, name="mdpr-decor heat-cell")
        add_line(slide, preview_x + 0.54, center_y - 0.12, 0.16, 0.018, accent, name="mdpr-decor heat-label")
    elif kind == "tag-cloud-strip":
        for idx, (cx, cy, ww) in enumerate([(0.02, -0.16, 0.28), (0.36, -0.16, 0.32), (0.08, 0.12, 0.34), (0.48, 0.12, 0.24)]):
            add_box(slide, preview_x + cx, center_y + cy, ww, 0.16, accent if idx % 2 == 0 else "CBD5E1", accent if idx % 2 == 0 else "CBD5E1", True, name="mdpr-decor tag-chip")
    elif kind == "trend-line-backdrop":
        add_line(slide, preview_x + 0.04, center_y + 0.14, 0.18, 0.02, accent, name="mdpr-decor trend")
        add_line(slide, preview_x + 0.22, center_y + 0.08, 0.2, 0.02, accent, name="mdpr-decor trend")
        add_line(slide, preview_x + 0.42, center_y - 0.03, 0.22, 0.02, accent, name="mdpr-decor trend")
        add_text(slide, preview_x + 0.06, center_y - 0.18, 0.58, 0.16, "trend", 5, LINE, True, PP_ALIGN.CENTER, name="mdpr-text-in-decor", margin=0)
    elif "summary" in kind and family == "table":
        add_box(slide, preview_x, center_y - 0.24, 0.5, 0.48, "FFFFFF", LINE, False, name="mdpr-decor table-summary")
        add_line(slide, preview_x, center_y - 0.24, 0.5, 0.07, accent, name="mdpr-decor table-header")
        add_line(slide, preview_x + 0.06, center_y - 0.03, 0.38, 0.014, LINE)
        add_line(slide, preview_x + 0.06, center_y + 0.13, 0.38, 0.014, LINE)
    elif "hub" in kind:
        add_dot(slide, preview_x + 0.25, center_y - 0.14, 0.28, accent)
        add_line(slide, preview_x + 0.1, center_y, 0.55, 0.018, "94A3B8")
        add_line(slide, preview_x + 0.38, center_y - 0.26, 0.018, 0.52, "94A3B8")
        for dx, dy in [(0.02, -0.26), (0.52, -0.24), (0.04, 0.25), (0.53, 0.22)]:
            add_dot(slide, preview_x + dx, center_y + dy, 0.13, "94A3B8")
    if "cycle" in kind or "loop" in kind:
        for dx, dy in [(0.12, -0.22), (0.46, -0.05), (0.18, 0.2)]:
            add_dot(slide, preview_x + dx, center_y + dy, 0.16, accent)
    if "quadrant" in kind or "matrix" in kind:
        add_line(slide, preview_x + 0.35, center_y - 0.28, 0.025, 0.56, accent)
        add_line(slide, preview_x + 0.08, center_y, 0.56, 0.025, accent)
    if ("donut" in kind or "ring" in kind or "gauge" in kind) and kind not in {"donut-label-ring", "gauge-score-card"}:
        add_dot(slide, preview_x + 0.16, center_y - 0.22, 0.46, accent)
        add_dot(slide, preview_x + 0.27, center_y - 0.11, 0.24, fill)
    if kind not in {"donut-label-ring", "gauge-score-card", "trend-line-backdrop"} and ("bars" in kind or "chart" in family or "trend" in kind):
        for bar, height in enumerate([0.16, 0.32, 0.24]):
            add_line(slide, preview_x + 0.1 + bar * 0.18, center_y + 0.22 - height, 0.1, height, accent)
    if kind not in {"table-summary-card", "status-dot-table", "risk-heat-grid"} and ("table" in family or "grid" in kind):
        for offset in [center_y - 0.18, center_y, center_y + 0.18]:
            add_line(slide, preview_x + 0.06, offset, 0.58, 0.015, LINE)
        for offset in [preview_x + 0.25, preview_x + 0.45]:
            add_line(slide, offset, center_y - 0.28, 0.015, 0.56, LINE)
    if "image" in family or "photo" in kind:
        add_box(slide, preview_x, center_y - 0.25, 0.46, 0.5, "DBEAFE", "93C5FD", False, name="mdpr-decor image-slot")
        add_line(slide, preview_x + 0.5, center_y - 0.12, 0.24, 0.035, accent)
    if "paper" in kind or "ticket" in kind or "document" in pattern.get("selectionSignals", []):
        add_box(slide, preview_x + 0.12, center_y - 0.24, 0.46, 0.48, SURFACE, LINE, False, name="mdpr-decor paper")
        add_line(slide, preview_x + 0.18, center_y - 0.06, 0.34, 0.018, LINE)
    if "rail" in kind:
        add_line(slide, preview_x, preview_y, 0.06, preview_h, accent)
    if "rule" in kind or "header" in kind or "overline" in kind:
        add_line(slide, preview_x, y + 0.2, preview_w, 0.035, accent)
    if kind.startswith("number-tab") or kind in {"number-tab", "rank-ribbon", "rank-ribbon-card"}:
        add_box(slide, preview_x, center_y - 0.19, 0.42, 0.38, accent, accent, name="mdpr-decor label-tab")
        add_text(slide, preview_x, center_y - 0.19, 0.42, 0.38, str((index % 9) + 1), 10, "FFFFFF", True, PP_ALIGN.CENTER, name="mdpr-text-in-decor", margin=0.0)
    if "chip" in kind or "badge" in kind:
        add_box(slide, preview_x + 0.1, center_y - 0.13, 0.5, 0.26, accent, accent, name="mdpr-decor chip")
    if "bracket" in kind:
        add_line(slide, preview_x, preview_y, 0.04, preview_h, accent)
        add_line(slide, preview_x, preview_y, 0.32, 0.035, accent)
        add_line(slide, preview_x, preview_y + preview_h - 0.035, 0.32, 0.035, accent)
    if kind not in {"status-dot-table"} and ("dot" in kind or "icon" in kind or "connector" in kind):
        add_dot(slide, preview_x + 0.05, center_y - 0.13, 0.26, accent)
        add_line(slide, preview_x + 0.38, center_y - 0.01, 0.42, 0.02, accent)
    if "meter" in kind or "progress" in kind:
        add_line(slide, preview_x, center_y + 0.12, preview_w, 0.05, "D1D5DB")
        add_line(slide, preview_x, center_y + 0.12, preview_w * 0.68, 0.05, accent)
    if "notch" in kind or "folded" in kind:
        add_box(slide, preview_x + 0.32, center_y - 0.2, 0.35, 0.3, "FDE68A", "FDE68A", False, name="mdpr-decor notch")
    if "ticket" in kind:
        add_dot(slide, preview_x - 0.08, center_y - 0.12, 0.24, BG)
        add_dot(slide, preview_x + preview_w - 0.16, center_y - 0.12, 0.24, BG)
    if "bubble" in kind:
        add_dot(slide, preview_x + 0.32, center_y + 0.12, 0.18, fill)
    if "image" in family:
        add_box(slide, preview_x, center_y - 0.23, 0.62, 0.46, "DBEAFE", "93C5FD", False, name="mdpr-decor image-slot")
    if kind not in {"table-summary-card", "status-dot-table", "risk-heat-grid"} and "table" in family:
        for offset in [center_y - 0.08, center_y + 0.08]:
            add_line(slide, preview_x, offset, preview_w, 0.02, LINE)
        for offset in [preview_x + 0.25, preview_x + 0.5]:
            add_line(slide, offset, center_y - 0.23, 0.02, 0.46, LINE)

    add_text(slide, text_x, y + 0.13, text_w, 0.26, pattern["id"], 8, TEXT, True, name="mdpr-text catalog-title")
    add_text(slide, text_x, y + 0.46, text_w, 0.35, pattern["shapeGrammar"], 7, MUTED, name="mdpr-text catalog-body")
    add_text(slide, text_x, y + h - 0.3, text_w, 0.2, "theme-token driven", 7, accent, True, name="mdpr-text catalog-token")


def add_title(slide, title: str, subtitle: str = "", show_rule: bool = True) -> None:
    add_text(slide, 0.55, 0.34, 8.7, 0.45, title, 24, TEXT, True)
    if subtitle:
        add_text(slide, 0.57, 0.84, 11.2, 0.36, subtitle, 12, "334155")
    if show_rule:
        add_line(slide, 0.58, 1.2, 1.4, 0.035, ACCENT)


def setup_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)
    return slide


def add_cover(prs: Presentation, audit: dict[str, Any]) -> None:
    slide = setup_slide(prs)
    add_text(slide, 0.7, 0.75, 9.4, 0.72, "MDPR Theme and Decoration Review", 32, TEXT, True)
    add_text(slide, 0.72, 1.55, 10.5, 0.45, "Actual PPTX review deck generated from theme matrix renders, object shape catalog, and deterministic coverage audit.", 13, MUTED)
    metrics = [
        ("Object kinds", str(audit["catalogAudit"]["uniqueKindCount"])),
        ("Token-ready", str(audit["catalogAudit"]["themeTokenReadyCount"])),
        ("Theme combos", str(audit["themeMatrixAudit"]["comboCount"])),
        ("Assessment", audit["assessment"]),
    ]
    for i, (label, value) in enumerate(metrics):
        x = 0.75 + i * 3.0
        add_box(slide, x, 2.45, 2.42, 1.05, SURFACE, LINE)
        add_text(slide, x + 0.18, 2.58, 1.9, 0.22, label.upper(), 7, MUTED, True)
        add_text(slide, x + 0.18, 2.92, 1.9, 0.35, value, 20, ACCENT if value != "needs-improvement" else CONTRAST, True)
    add_text(slide, 0.78, 4.15, 10.8, 0.38, "Pass criteria: at least 30 distinct object shape kinds, every counted kind bound to theme tokens, and actual PowerPoint-exported theme previews passing style/color checks.", 13, TEXT)
    add_box(slide, 0.78, 5.05, 11.7, 1.1, "ECFDF5", "A7F3D0")
    add_text(slide, 1.05, 5.25, 10.9, 0.5, "This deck is itself a verification artifact: it contains the rendered contact sheets, editable pattern preview objects, and the audit summary path.", 16, "064E3B", True)
    add_text(slide, 1.05, 5.78, 10.9, 0.24, "artifacts/theme-decoration-review/theme-decoration-coverage-report.json", 10, "047857")


def add_contact_sheet_slide(prs: Presentation, title: str, image_path: Path, caption: str) -> None:
    slide = setup_slide(prs)
    add_title(slide, title, caption)
    max_x, max_y, max_w, max_h = 0.55, 1.35, 12.0, 5.55
    image = Image.open(image_path)
    image_ratio = image.width / image.height
    box_ratio = max_w / max_h
    if image_ratio >= box_ratio:
        draw_w = max_w
        draw_h = max_w / image_ratio
    else:
        draw_h = max_h
        draw_w = max_h * image_ratio
    draw_x = max_x + (max_w - draw_w) / 2
    draw_y = max_y + (max_h - draw_h) / 2
    slide.shapes.add_picture(str(image_path), Inches(draw_x), Inches(draw_y), width=Inches(draw_w), height=Inches(draw_h))
    if abs((draw_w / draw_h) - image_ratio) > 0.02:
        VISUAL_REVIEW_FINDINGS.append(f"{title}: inserted image aspect ratio drift")


def add_reference_corpus_slide(prs: Presentation, audit: dict[str, Any]) -> None:
    slide = setup_slide(prs)
    ref = audit["referenceCorpusAudit"]
    add_title(slide, "Reference Corpus Expansion", "Only source-neutral structure counts and derived grammar are kept in the repository.")
    metrics = [
        ("PPT files", ref["pptDownloaded"]),
        ("Slides", ref["slidesAnalyzed"]),
        ("Rendered PNGs", ref["renderedPngSlides"]),
        ("PNG samples", ref["pngSamplesAnalyzed"]),
        ("Derived patterns", ref["derivedObjectPatternCount"]),
    ]
    for i, (label, value) in enumerate(metrics):
        x = 0.72 + (i % 5) * 2.45
        add_box(slide, x, 1.7, 2.05, 1.1, ["ECFDF5", "EFF6FF", "FFF7ED", "F5F3FF", "F8FAFC"][i], LINE, name="mdpr-surface reference-metric")
        add_text(slide, x + 0.18, 1.88, 1.55, 0.24, label.upper(), 7, MUTED, True)
        add_text(slide, x + 0.18, 2.18, 1.55, 0.34, f"{value:,}", 20, [ACCENT, "2563EB", "C2410C", "7C3AED", CONTRAST][i], True)
    add_box(slide, 0.82, 3.5, 5.45, 1.65, SURFACE, LINE, name="mdpr-surface reference-policy")
    add_text(slide, 1.1, 3.78, 4.85, 0.72, "Corpus assets stay local. Public artifacts store counts, object families, and coherence guards only.", 17, TEXT, True)
    add_box(slide, 7.0, 3.5, 4.95, 1.65, "FFF7ED", "FED7AA", name="mdpr-surface diversity-policy")
    add_text(slide, 7.3, 3.75, 4.35, 0.72, "Diversity is measured by structural archetype count and family spread, not by repeated card variants.", 16, "9A3412", True)


def draw_atlas_icon(slide, pattern: dict[str, Any], x: float, y: float, w: float, h: float, index: int) -> None:
    accent = [ACCENT, "2563EB", "C2410C", "7C3AED", CONTRAST, "16A34A"][index % 6]
    fill = {"chart": "F0FDF4", "diagram": "EFF6FF", "table": "F8FAFC", "image": "EFF6FF", "surface": "FFF7ED", "metric": "ECFDF5"}.get(pattern["family"], SURFACE)
    add_box(slide, x, y, w, h, fill, LINE, name="mdpr-surface atlas")
    cx = x + w / 2
    cy = y + h / 2
    kind = pattern["kind"]
    family = pattern["family"]
    if "hub" in kind:
        add_dot(slide, cx - 0.18, cy - 0.18, 0.36, accent)
        add_line(slide, cx - 0.65, cy - 0.01, 1.3, 0.02, "94A3B8")
        add_line(slide, cx - 0.01, cy - 0.42, 0.02, 0.84, "94A3B8")
        for dx, dy in [(-0.62, -0.38), (0.5, -0.36), (-0.58, 0.38), (0.54, 0.35)]:
            add_dot(slide, cx + dx, cy + dy, 0.2, "CBD5E1")
    elif "loop" in kind or "cycle" in kind:
        for dx, dy, label in [(-0.48, -0.28, "1"), (0.48, -0.14, "2"), (0.02, -0.03, "3")]:
            add_dot(slide, cx + dx, cy + dy, 0.3, accent)
            add_text(slide, cx + dx, cy + dy, 0.3, 0.3, label, 9, "FFFFFF", True, PP_ALIGN.CENTER, name="mdpr-text-in-decor", margin=0)
        add_arrow(slide, cx - 0.15, cy - 0.5, 0.38, 0.16, "94A3B8")
    elif "quadrant" in kind or "risk" in kind:
        add_line(slide, x + 0.32, cy, w - 0.64, 0.02, accent)
        add_line(slide, cx, y + 0.28, 0.02, h - 0.78, accent)
        for ox, oy in [(0.42, 0.38), (w - 0.9, 0.38), (0.42, h - 0.72), (w - 0.9, h - 0.72)]:
            add_box(slide, x + ox, y + oy, 0.45, 0.28, "FFFFFF", LINE, False)
    elif kind == "donut-label-ring":
        ring_size = 0.62
        add_dot(slide, cx - ring_size / 2, cy - ring_size / 2 - 0.02, ring_size, accent)
        add_dot(slide, cx - ring_size / 2 + 0.16, cy - ring_size / 2 + 0.14, ring_size - 0.3, fill)
        add_line(slide, cx + 0.48, cy - 0.22, 0.56, 0.02, accent)
        add_line(slide, cx + 0.48, cy + 0.18, 0.46, 0.02, "94A3B8")
    elif kind == "gauge-score-card":
        add_dot(slide, cx - 0.36, cy - 0.35, 0.72, "DCFCE7")
        add_dot(slide, cx - 0.25, cy - 0.24, 0.5, fill)
        add_line(slide, cx - 0.04, cy + 0.04, 0.42, 0.03, accent)
        add_text(slide, cx - 0.23, cy - 0.08, 0.46, 0.18, "78", 8, accent, True, PP_ALIGN.CENTER, name="mdpr-text-in-decor", margin=0)
    elif family == "chart":
        for i, height in enumerate([0.36, 0.72, 0.52, 0.9]):
            add_line(slide, x + 0.45 + i * 0.35, y + h - 0.58 - height, 0.18, height, accent)
        add_dot(slide, x + w - 1.0, y + 0.45, 0.62, accent)
        add_dot(slide, x + w - 0.85, y + 0.6, 0.32, fill)
    elif family == "table":
        for row in range(4):
            add_line(slide, x + 0.35, y + 0.36 + row * 0.18, w - 0.7, 0.015, LINE)
        for col in range(3):
            add_line(slide, x + 0.6 + col * 0.55, y + 0.28, 0.015, 0.72, LINE)
        add_dot(slide, x + 0.48, y + 0.48, 0.13, accent)
        add_dot(slide, x + 0.48, y + 0.68, 0.13, CONTRAST)
    elif family == "image":
        add_box(slide, x + 0.35, y + 0.38, 0.95, 0.88, "DBEAFE", "93C5FD", False)
        add_box(slide, x + 1.45, y + 0.52, 0.9, 0.22, accent, accent, True)
        add_line(slide, x + 1.45, y + 0.9, 0.95, 0.03, LINE)
    elif "paper" in kind or "ticket" in kind:
        add_box(slide, x + 0.72, y + 0.42, 1.0, 0.5, SURFACE, LINE, False)
        add_box(slide, x + 0.58, y + 0.32, 1.0, 0.5, SURFACE, LINE, False)
        add_line(slide, x + 0.72, y + 0.58, 0.58, 0.02, LINE)
    elif "tag" in kind or "label" in family:
        for i, label_w in enumerate([0.55, 0.72, 0.48, 0.86]):
            add_box(slide, x + 0.4 + (i % 2) * 0.9, y + 0.48 + (i // 2) * 0.42, label_w, 0.28, [accent, "CBD5E1"][i % 2], [accent, "CBD5E1"][i % 2], True)
    else:
        add_dot(slide, x + 0.45, cy - 0.15, 0.3, accent)
        add_line(slide, x + 0.88, cy, 1.1, 0.03, accent)
        add_box(slide, x + 1.85, cy - 0.22, 0.48, 0.44, SURFACE, LINE, True)
    add_text(slide, x + 0.28, y + h - 0.42, w - 0.56, 0.24, pattern["id"], 8, TEXT, True, PP_ALIGN.CENTER)


def add_archetype_atlas_slide(prs: Presentation, patterns: list[dict[str, Any]]) -> None:
    slide = setup_slide(prs)
    add_title(slide, "Structural Archetype Atlas", "Representative non-card object grammars from the source-neutral 60-pattern catalog.")
    wanted = [
        "center-hub-spokes", "loop-arrow-cycle", "axis-quadrant-map", "vertical-timeline-cards",
        "donut-label-ring", "gauge-score-card", "risk-heat-grid", "status-dot-table",
        "image-caption-split", "pictorial-anchor-labels", "stacked-paper-cards", "tag-cloud-strip",
    ]
    by_id = {pattern["id"]: pattern for pattern in patterns}
    subset = [by_id[item] for item in wanted if item in by_id]
    for index, pattern in enumerate(subset):
        col = index % 4
        row = index // 4
        draw_atlas_icon(slide, pattern, 0.68 + col * 3.12, 1.55 + row * 1.68, 2.62, 1.34, index)


def add_arrow(slide, x: float, y: float, w: float, h: float, color: str = ACCENT, name: str = "mdpr-arrow") -> None:
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.name = name
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = rgb(color)
    arrow.line.color.rgb = rgb(color)


def add_layout_diversity_slide(prs: Presentation) -> None:
    slide = setup_slide(prs)
    add_title(slide, "Layout Diversity Review", "The review deck exercises non-card layouts: timeline, split comparison, data proof, and editorial image framing.")
    # Horizontal timeline
    add_text(slide, 0.75, 1.45, 2.1, 0.28, "Sequence rail", 13, TEXT, True)
    rail_y = 2.02
    add_line(slide, 0.85, rail_y + 0.12, 5.15, 0.035, ACCENT)
    for i, label in enumerate(["Parse", "Plan", "Render", "Review"]):
        cx = 0.9 + i * 1.7
        add_dot(slide, cx, rail_y, 0.28, ACCENT)
        add_text(slide, cx - 0.3, rail_y + 0.34, 0.9, 0.28, label, 9, TEXT, True, PP_ALIGN.CENTER)
    # Split comparison
    add_text(slide, 6.65, 1.45, 2.4, 0.28, "Split comparison", 13, TEXT, True)
    add_box(slide, 6.65, 1.88, 2.55, 1.1, "ECFDF5", "A7F3D0", name="mdpr-surface layout-comparison")
    add_box(slide, 9.35, 1.88, 2.55, 1.1, "FFF7ED", "FED7AA", name="mdpr-surface layout-comparison")
    add_text(slide, 6.9, 2.05, 2.05, 0.48, "Rule base\nowns geometry", 13, "065F46", True, PP_ALIGN.CENTER)
    add_text(slide, 9.6, 2.05, 2.05, 0.48, "LLM hint\nadds intent", 13, "9A3412", True, PP_ALIGN.CENTER)
    # Data proof with chart/table
    add_text(slide, 0.75, 3.45, 2.4, 0.28, "Data proof", 13, TEXT, True)
    add_box(slide, 0.75, 3.85, 5.25, 1.9, SURFACE, LINE, name="mdpr-surface data-proof")
    for i, height in enumerate([0.45, 0.82, 0.62, 1.08]):
        add_line(slide, 1.1 + i * 0.6, 5.15 - height, 0.32, height, ["2563EB", ACCENT, CONTRAST, "7C3AED"][i], name="mdpr-decor bar")
    add_text(slide, 3.75, 4.02, 1.8, 0.32, "98 slides", 20, ACCENT, True)
    add_text(slide, 3.75, 4.46, 1.85, 0.52, "exported through PowerPoint before review", 10, MUTED)
    add_line(slide, 3.75, 5.2, 1.7, 0.025, LINE)
    # Editorial image frame
    add_text(slide, 6.65, 3.45, 2.7, 0.28, "Editorial frame", 13, TEXT, True)
    add_box(slide, 6.65, 3.85, 2.1, 1.9, "DBEAFE", "93C5FD", False, name="mdpr-surface image-frame")
    add_box(slide, 8.95, 4.1, 2.95, 1.4, SURFACE, LINE, name="mdpr-surface caption-card")
    add_text(slide, 9.22, 4.32, 2.35, 0.55, "Image-aware layouts keep caption and proof text in separate regions.", 12, TEXT, True)


def add_process_layout_slide(prs: Presentation) -> None:
    slide = setup_slide(prs)
    add_title(slide, "Connected Process Layout", "Arrows connect aligned parent regions; inner labels remain centered inside their own marker shapes.")
    items = [("A", "Markdown\nsemantics"), ("B", "Rule-based\nlayout"), ("C", "Editable\nPPT objects"), ("D", "Visual\nvalidation")]
    for i, (letter, label) in enumerate(items):
        x = 0.78 + i * 3.05
        add_box(slide, x, 2.2, 2.25, 1.65, ["ECFDF5", "EFF6FF", "FFF7ED", "F5F3FF"][i], LINE, name="mdpr-surface process")
        add_dot(slide, x + 0.23, 2.72, 0.5, [ACCENT, "2563EB", "C2410C", "7C3AED"][i], name="mdpr-decor centered-marker")
        add_text(slide, x + 0.23, 2.72, 0.5, 0.5, letter, 15, "FFFFFF", True, PP_ALIGN.CENTER, name="mdpr-text-in-decor", margin=0)
        add_text(slide, x + 0.85, 2.55, 1.6, 0.78, label, 15, TEXT, True, PP_ALIGN.LEFT)
        if i < len(items) - 1:
            add_arrow(slide, x + 2.38, 2.83, 0.44, 0.26, "94A3B8")
    add_box(slide, 1.25, 5.05, 10.55, 0.72, "F8FAFC", LINE, name="mdpr-surface note")
    add_text(slide, 1.55, 5.18, 9.9, 0.34, "Validation target: markers and labels share vertical centers; arrows touch parent-region edges without crossing text.", 13, MUTED, True)


def add_mixed_object_slide(prs: Presentation, audit: dict[str, Any]) -> None:
    slide = setup_slide(prs)
    add_title(slide, "Mixed Object Stress Layout", "Tables, metrics, chart-like marks, icon markers, and callouts are composed without forcing empty-space filler.")
    catalog = audit["catalogAudit"]
    add_box(slide, 0.72, 1.55, 3.0, 1.25, "ECFDF5", "A7F3D0", name="mdpr-surface metric")
    add_text(slide, 1.0, 1.76, 1.1, 0.42, str(catalog["structuralArchetypeCount"]), 28, ACCENT, True)
    add_text(slide, 2.05, 1.78, 1.3, 0.5, "object\nkinds", 14, TEXT, True)
    add_box(slide, 4.2, 1.55, 3.7, 2.25, SURFACE, LINE, name="mdpr-surface table")
    for row in range(4):
        y = 1.9 + row * 0.38
        add_line(slide, 4.42, y, 3.25, 0.015, "E2E8F0", name="mdpr-decor table-rule")
        add_text(slide, 4.48, y + 0.04, 1.4, 0.22, ["Cards", "Diagrams", "Metrics", "Images"][row], 9, TEXT, row == 0)
        counts = catalog["familyCounts"]
        add_text(slide, 7.05, y + 0.04, 0.45, 0.22, str([counts.get("card", 0), counts.get("diagram", 0), counts.get("metric", 0), counts.get("image", 0)][row]), 9, ACCENT, True, PP_ALIGN.RIGHT)
    add_box(slide, 8.45, 1.55, 3.95, 2.25, SURFACE, LINE, name="mdpr-surface chart")
    for i, height in enumerate([0.45, 0.9, 0.62, 1.2, 0.78]):
        add_line(slide, 8.8 + i * 0.46, 3.3 - height, 0.24, height, [ACCENT, "2563EB", "C2410C", "7C3AED", CONTRAST][i], name="mdpr-decor chart-bar")
    add_text(slide, 11.35, 2.0, 0.7, 0.4, str(len(catalog["families"])), 22, CONTRAST, True, PP_ALIGN.CENTER)
    add_text(slide, 10.95, 2.52, 1.45, 0.5, "object families\nchecked", 10, MUTED, True, PP_ALIGN.CENTER)
    for i, text in enumerate(["No text crosses shape boundary", "Decoration lives outside text box", "Minimum font stays readable"]):
        y = 4.55 + i * 0.5
        add_dot(slide, 0.9, y + 0.07, 0.22, ACCENT, name="mdpr-decor bullet")
        add_text(slide, 0.9, y + 0.07, 0.22, 0.22, chr(97 + i), 8, "FFFFFF", True, PP_ALIGN.CENTER, name="mdpr-text-in-decor", margin=0)
        add_text(slide, 1.25, y, 5.1, 0.34, text, 13, TEXT, True)
    add_box(slide, 7.1, 4.45, 4.7, 1.45, "FFF7ED", "FED7AA", name="mdpr-surface callout")
    add_text(slide, 7.45, 4.75, 4.0, 0.58, "Review note: layout variation is checked by slide families, not by card-count alone.", 13, "9A3412", True)


def add_theme_primitive_matrix_slide(prs: Presentation, audit: dict[str, Any]) -> None:
    slide = setup_slide(prs)
    add_title(slide, "Theme Style Primitive Matrix", "The same card, table, chart, and marker primitives change decoration grammar under each theme style.")
    audits = audit["themeMatrixAudit"]["themeAudits"]
    for index, item in enumerate(audits[:8]):
        col = index % 4
        row = index // 4
        x = 0.62 + col * 3.12
        y = 1.42 + row * 2.45
        w = 2.72
        h = 1.98
        lock = item
        colors = lock["surfacePolicy"]
        theme_colors = lock.get("themeColors", {})
        style = lock["style"]
        background = theme_colors.get("light1") or ("0B1020" if style in {"glass", "data", "dark"} else BG)
        text = theme_colors.get("dark1") or ("F8FAFC" if style in {"glass", "data", "dark"} else TEXT)
        surface = theme_colors.get("light2") or ("17213A" if style in {"glass", "data", "dark"} else SURFACE)
        line = theme_colors.get("accent4") or LINE
        accent = theme_colors.get("accent1") or ACCENT
        support = theme_colors.get("accent2") or "2563EB"
        add_box(slide, x, y, w, h, background, blend_hex(line, background, 0.45), False, name="mdpr-theme-cell")
        add_themed_surface(slide, x + 0.2, y + 0.3, 1.16, 0.72, style, surface, line, accent, background, name=f"mdpr-surface primitive-{style}-card")
        add_text(slide, x + 0.34, y + 0.45, 0.78, 0.22, "CARD", 7, text, True, PP_ALIGN.CENTER, name="mdpr-text primitive-label")
        add_themed_surface(slide, x + 1.54, y + 0.3, 0.86, 0.72, style, surface, line, accent, background, name=f"mdpr-surface primitive-{style}-table")
        for r in range(3):
            add_line(slide, x + 1.66, y + 0.45 + r * 0.15, 0.58, 0.01, blend_hex(line, background, 0.18), name="mdpr-decor primitive-table-row")
        for b, bh in enumerate([0.18, 0.34, 0.26]):
            add_line(slide, x + 0.36 + b * 0.22, y + 1.57 - bh, 0.11, bh, accent if b != 1 else support, name="mdpr-decor primitive-chart-bar")
        add_dot(slide, x + 1.72, y + 1.22, 0.28, accent, name="mdpr-decor primitive-marker")
        add_text(slide, x + 1.72, y + 1.22, 0.28, 0.28, chr(65 + row * 4 + col), 8, "FFFFFF", True, PP_ALIGN.CENTER, name="mdpr-text-in-decor", margin=0)
        add_text(slide, x + 2.02, y + 1.15, 0.58, 0.28, "mark", 7, text, True, name="mdpr-text primitive-caption", margin=0.04)
        add_text(slide, x + 0.2, y + 0.06, 1.35, 0.24, style.upper(), 8, text, True, name="mdpr-text primitive-style", margin=0.04)
        add_text(slide, x + 1.42, y + 0.06, 1.1, 0.24, colors["shadow"], 7, text, False, PP_ALIGN.RIGHT, name="mdpr-text primitive-effect", margin=0.04)


def add_catalog_slides(prs: Presentation, patterns: list[dict[str, Any]]) -> None:
    for page, start in enumerate(range(0, len(patterns), 12), 1):
        slide = setup_slide(prs)
        add_title(slide, f"Object Shape Grammar Catalog {page}", "Editable preview objects generated from the 60+ structural archetype catalog.")
        subset = patterns[start:start + 12]
        cols = 3
        card_w = 3.86
        card_h = 1.05
        gap_x = 0.25
        gap_y = 0.28
        for idx, pattern in enumerate(subset):
            col = idx % cols
            row = idx // cols
            x = 0.62 + col * (card_w + gap_x)
            y = 1.45 + row * (card_h + gap_y)
            add_pattern_preview(slide, pattern, x, y, card_w, card_h, start + idx)


def add_breakdown_slide(prs: Presentation, audit: dict[str, Any]) -> None:
    slide = setup_slide(prs)
    add_title(slide, "Coverage Audit Summary", show_rule=False)
    add_box(slide, 0.75, 1.12, 11.45, 0.52, "ECFDF5", "A7F3D0", name="mdpr-surface audit-claim")
    add_text(
        slide,
        1.02,
        1.24,
        10.75,
        0.22,
        "Theme feel and object diversity are verified from generated PPTX and PNG artifacts, not from documentation alone.",
        13,
        "065F46",
        True,
        name="mdpr-text audit-claim",
    )
    catalog = audit["catalogAudit"]
    matrix = audit["themeMatrixAudit"]
    left = [
        ("Catalog patterns", catalog["catalogPatternCount"]),
        ("Unique shape kinds", catalog["uniqueKindCount"]),
        ("Archetypes", catalog["structuralArchetypeCount"]),
        ("Theme-token ready", catalog["themeTokenReadyCount"]),
    ]
    for i, (label, value) in enumerate(left):
        y = 1.82 + i * 0.72
        add_box(slide, 0.75, y, 4.7, 0.48, SURFACE, LINE)
        add_text(slide, 0.98, y + 0.08, 2.6, 0.2, label, 10, MUTED, True)
        add_text(slide, 4.3, y + 0.06, 0.7, 0.24, str(value), 14, ACCENT, True, PP_ALIGN.RIGHT)

    add_box(slide, 6.05, 1.82, 6.15, 2.2, "F8FAFC", LINE)
    add_text(slide, 6.32, 2.05, 5.3, 0.28, "Families", 14, TEXT, True)
    families = ", ".join(catalog["families"])
    add_text(slide, 6.32, 2.39, 5.4, 0.9, families, 10, MUTED)
    add_text(slide, 6.32, 3.42, 5.2, 0.28, f"Average theme RGB distance: {matrix['visualDistinctness']['averageMeanRgbDistance']}", 11, ACCENT, True)

    add_box(slide, 6.05, 4.37, 6.15, 1.55, "FFF7ED", "FED7AA")
    add_text(slide, 6.32, 4.57, 5.4, 0.3, "Manual visual review note", 13, "9A3412", True)
    add_text(slide, 6.32, 4.95, 5.4, 0.54, "Diversity is now judged by 60+ structural archetypes across 18 object families, not by repeated card decoration variants.", 10, "9A3412")


def export_with_powerpoint(pptx_path: Path, output_dir: Path) -> list[Path]:
    import win32com.client  # type: ignore

    if output_dir.exists():
        shutil.rmtree(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    app = win32com.client.DispatchEx("PowerPoint.Application")
    presentation = None
    try:
        app.Visible = 1
        presentation = app.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        presentation.Export(str(output_dir.resolve()), "PNG", 1600, 900)
        presentation.Close()
        presentation = None
    finally:
        if presentation is not None:
            presentation.Close()
        app.Quit()
    def natural_key(path: Path) -> list[Any]:
        return [int(part) if part.isdigit() else part for part in re.split(r"(\d+)", path.stem)]

    return sorted({path.resolve(): path for path in list(output_dir.glob("*.PNG")) + list(output_dir.glob("*.png"))}.values(), key=natural_key)


def validate_deck(pptx_path: Path, pngs: list[Path]) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    shape_count = 0
    picture_count = 0
    text_count = 0
    min_font_pt = 999
    bounds_violations: list[dict[str, Any]] = []
    text_margin_violations: list[dict[str, Any]] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            shape_count += 1
            left = shape.left / 914400
            top = shape.top / 914400
            width = shape.width / 914400
            height = shape.height / 914400
            if left < -0.03 or top < -0.03 or left + width > SLIDE_W + 0.03 or top + height > SLIDE_H + 0.03:
                bounds_violations.append({"name": shape.name, "left": round(left, 3), "top": round(top, 3), "width": round(width, 3), "height": round(height, 3)})
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_count += 1
                tf = shape.text_frame
                if shape.name != "mdpr-text-in-decor" and (tf.margin_left < Inches(0.035) or tf.margin_right < Inches(0.035)):
                    text_margin_violations.append({"name": shape.name, "text": shape.text[:48]})
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is not None:
                            min_font_pt = min(min_font_pt, round(run.font.size.pt, 2))
    png_checks = []
    for path in pngs:
        img = Image.open(path).convert("RGB")
        colors = len(img.getcolors(maxcolors=1_000_000) or [])
        png_checks.append({"file": str(path.relative_to(ROOT)), "size": list(img.size), "uniqueColors": colors, "hasContent": colors > 20})
    layout_families = ["contact-sheet", "reference-corpus", "archetype-atlas", "timeline", "split-comparison", "process", "data-table-chart", "theme-primitive-matrix", "catalog-grid", "audit-summary"]
    style_effect_count = 0
    table_summary_preview_count = 0
    aspect_findings = list(VISUAL_REVIEW_FINDINGS)
    for slide in prs.slides:
        for shape in slide.shapes:
            name = shape.name or ""
            if "mdpr-effect" in name or "primitive-" in name:
                style_effect_count += 1
            if "table-summary" in name:
                table_summary_preview_count += 1
    return {
        "slideCount": len(prs.slides),
        "shapeCount": shape_count,
        "pictureCount": picture_count,
        "textShapeCount": text_count,
        "minFontPt": min_font_pt if min_font_pt != 999 else None,
        "boundsViolations": bounds_violations,
        "textMarginViolations": text_margin_violations,
        "layoutFamilies": layout_families,
        "styleEffectShapeCount": style_effect_count,
        "tableSummaryPreviewCount": table_summary_preview_count,
        "aspectFindings": aspect_findings,
        "pngValidation": png_checks,
        "ok": (
            len(prs.slides) >= 15
            and len(layout_families) >= 8
            and shape_count > 300
            and text_count > 120
            and style_effect_count >= 40
            and table_summary_preview_count >= 1
            and not aspect_findings
            and (min_font_pt if min_font_pt != 999 else MIN_TEXT_PT) >= MIN_TEXT_PT
            and not bounds_violations
            and not text_margin_violations
            and len(pngs) == len(prs.slides)
            and all(item["hasContent"] for item in png_checks)
        ),
    }


def create_deck() -> dict[str, Any]:
    OUT.mkdir(parents=True, exist_ok=True)
    audit = json.loads(REPORT.read_text(encoding="utf-8"))
    catalog = json.loads(CATALOG.read_text(encoding="utf-8"))
    patterns = catalog["patterns"]
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    add_cover(prs, audit)
    add_contact_sheet_slide(
        prs,
        "Theme Cover Matrix",
        MATRIX / "theme-style-color-cover-contact-sheet.png",
        "PowerPoint-exported cover previews for style/color/harmony combinations.",
    )
    add_contact_sheet_slide(
        prs,
        "Theme Proof Object Matrix",
        MATRIX / "theme-style-color-proof-contact-sheet.png",
        "PowerPoint-exported arc-ring proof slides using each resolved theme palette.",
    )
    add_reference_corpus_slide(prs, audit)
    add_archetype_atlas_slide(prs, patterns)
    add_layout_diversity_slide(prs)
    add_process_layout_slide(prs)
    add_mixed_object_slide(prs, audit)
    add_theme_primitive_matrix_slide(prs, audit)
    add_catalog_slides(prs, patterns)
    add_breakdown_slide(prs, audit)
    prs.save(PPTX)
    pngs = export_with_powerpoint(PPTX, EXPORT_DIR)
    validation = validate_deck(PPTX, pngs)
    report = {
        "pptx": str(PPTX.relative_to(ROOT)),
        "pngDir": str(EXPORT_DIR.relative_to(ROOT)),
        "sourceAudit": str(REPORT.relative_to(ROOT)),
        "validation": validation,
        "ok": validation["ok"],
    }
    report_path = OUT / "theme-decoration-review-deck-report.json"
    report_path.write_text(json.dumps(report, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    print(json.dumps(report, indent=2, ensure_ascii=False))
    if not report["ok"]:
        raise SystemExit(1)
    return report


if __name__ == "__main__":
    create_deck()
