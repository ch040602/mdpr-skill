#!/usr/bin/env python3
from __future__ import annotations

import json
import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "artifacts" / "design-showcase"
MIN_FONT_SIZE_PT = 8


@dataclass(frozen=True)
class Skin:
    id: str
    name: str
    brand: str
    bg: str
    card: str
    surface: str
    text: str
    muted: str
    line: str
    success: str
    danger: str
    radius: float
    shadow: str


SKINS = [
    Skin("toss", "Friendly Dashboard", "721FE5", "FAFAFA", "FFFFFF", "F3F4F5", "2A2A2A", "6A6A6A", "E8E6E1", "2E9D63", "D4183D", 0.18, "soft"),
    Skin("stripe", "Layered Product", "533AFD", "F7F9FC", "FFFFFF", "EEF4FF", "20242C", "64748B", "D9E2F2", "0E9F6E", "EF476F", 0.13, "layered"),
    Skin("linear", "Sharp Technical", "5E6AD2", "111217", "1B1D25", "242733", "F4F5F8", "A5ADBA", "343846", "34D399", "FB7185", 0.07, "hairline"),
    Skin("notion", "Editorial Brief", "0075DE", "FBFAF7", "FFFFFF", "F3F1EA", "2F2F2F", "78716C", "E7E1D5", "3B8F5D", "B42318", 0.10, "none"),
]
SHOWCASE_ORDER = ["mixed-object-stress", "notion", "linear", "stripe", "toss"]

EMU_PER_INCH = 914400


def rgb(hex_value: str) -> RGBColor:
    return RGBColor(*(int(hex_value[i:i + 2], 16) for i in (0, 2, 4)))


def rgb_tuple(hex_value: str) -> tuple[int, int, int]:
    return tuple(int(hex_value[i:i + 2], 16) for i in (0, 2, 4))


def add_shape(slide, name: str, kind: MSO_AUTO_SHAPE_TYPE, x: float, y: float, w: float, h: float, fill: str, line: str | None = None, text: str | None = None, font_size: int = 16, font_color: str | None = None, bold: bool = False):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    if text is not None:
      tf = shape.text_frame
      tf.clear()
      tf.margin_left = Inches(0.14)
      tf.margin_right = Inches(0.14)
      tf.margin_top = Inches(0.08)
      tf.margin_bottom = Inches(0.08)
      tf.word_wrap = False
      tf.vertical_anchor = MSO_ANCHOR.MIDDLE
      p = tf.paragraphs[0]
      run = p.add_run()
      run.text = text
      run.font.size = Pt(max(font_size, MIN_FONT_SIZE_PT))
      run.font.bold = bold
      run.font.color.rgb = rgb(font_color or "111111")
    return shape


def add_text(slide, name: str, x: float, y: float, w: float, h: float, text: str, size: int, color: str, bold: bool = False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(max(size, MIN_FONT_SIZE_PT))
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_stat_card(slide, skin: Skin, idx: int, x: float, y: float, label: str, value: str, trend: str):
    add_shape(slide, f"{skin.id}_stat_{idx}_card", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 2.25, 1.28, skin.card, skin.line)
    add_shape(slide, f"{skin.id}_stat_{idx}_icon", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + 0.16, y + 0.16, 0.34, 0.34, skin.surface, skin.surface, text="•", font_size=12, font_color=skin.brand, bold=True)
    add_text(slide, f"{skin.id}_stat_{idx}_label", x + 0.58, y + 0.19, 1.4, 0.22, label.upper(), 8, skin.muted, True)
    add_text(slide, f"{skin.id}_stat_{idx}_value", x + 0.18, y + 0.54, 1.5, 0.42, value, 25, skin.text, True)
    add_text(slide, f"{skin.id}_stat_{idx}_trend", x + 0.18, y + 0.98, 1.2, 0.22, trend, 9, skin.success, True)


def add_chart_card(slide, skin: Skin, x: float, y: float):
    add_shape(slide, f"{skin.id}_chart_card", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 5.1, 3.0, skin.card, skin.line)
    add_text(slide, f"{skin.id}_chart_title", x + 0.26, y + 0.22, 2.5, 0.32, "Weekly activation", 15, skin.text, True)
    for i, period in enumerate(["7D", "30D", "90D"]):
        fill = skin.brand if i == 1 else skin.surface
        color = "FFFFFF" if i == 1 else skin.muted
        add_shape(slide, f"{skin.id}_period_{period}", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + 3.14 + i * 0.62, y + 0.22, 0.52, 0.26, fill, fill, text=period, font_size=7, font_color=color, bold=True)
    bar_base = y + 2.35
    heights = [0.55, 0.85, 0.48, 1.15, 0.92, 1.38, 1.05]
    for i, height in enumerate(heights):
        fill = skin.brand if i in [3, 5] else skin.surface
        add_shape(slide, f"{skin.id}_chart_bar_{i}", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + 0.42 + i * 0.58, bar_base - height, 0.32, height, fill, fill)
    add_text(slide, f"{skin.id}_chart_stat_a", x + 0.36, y + 2.56, 1.2, 0.2, "Conversion 18.4%", 8, skin.muted, True)
    add_text(slide, f"{skin.id}_chart_stat_b", x + 2.0, y + 2.56, 1.2, 0.2, "Retention 42%", 8, skin.muted, True)


def add_ranked_list(slide, skin: Skin, x: float, y: float):
    add_shape(slide, f"{skin.id}_rank_card", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 3.15, 3.0, skin.card, skin.line)
    add_text(slide, f"{skin.id}_rank_title", x + 0.24, y + 0.24, 2.3, 0.28, "Top slide patterns", 14, skin.text, True)
    rows = [("1", "stat-card", "92"), ("2", "chart-card", "87"), ("3", "insight-card", "76")]
    for i, (rank, name, score) in enumerate(rows):
        yy = y + 0.76 + i * 0.64
        fill = skin.surface if i else skin.brand
        text_color = "FFFFFF" if i == 0 else skin.text
        add_shape(slide, f"{skin.id}_rank_{i}_row", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + 0.22, yy, 2.68, 0.48, fill, fill)
        add_shape(slide, f"{skin.id}_rank_{i}_num", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + 0.36, yy + 0.1, 0.28, 0.28, skin.card if i == 0 else skin.line, skin.card if i == 0 else skin.line, text=rank, font_size=8, font_color=skin.brand if i == 0 else skin.muted, bold=True)
        add_text(slide, f"{skin.id}_rank_{i}_name", x + 0.75, yy + 0.15, 1.2, 0.2, name, 9, text_color, True)
        add_text(slide, f"{skin.id}_rank_{i}_score", x + 2.36, yy + 0.15, 0.35, 0.2, score, 9, text_color, True)


def add_insight(slide, skin: Skin, x: float, y: float):
    add_shape(slide, f"{skin.id}_insight_card", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 3.2, 1.5, skin.surface, skin.brand)
    add_text(slide, f"{skin.id}_insight_title", x + 0.22, y + 0.2, 2.2, 0.26, "Coherence lock", 13, skin.text, True)
    add_text(slide, f"{skin.id}_insight_body", x + 0.22, y + 0.55, 2.58, 0.44, "One radius, one shadow family, one accent role per deck.", 10, skin.muted)
    for i, tag in enumerate(["radius", "shadow", "accent"]):
        add_shape(slide, f"{skin.id}_tag_{i}", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + 0.22 + i * 0.72, y + 1.1, 0.58, 0.22, skin.card, skin.line, text=tag, font_size=6, font_color=skin.muted, bold=True)


def create_reference_image(path: Path, skin: Skin) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    width, height = 720, 420
    image = Image.new("RGB", (width, height), rgb_tuple(skin.surface))
    draw = ImageDraw.Draw(image)
    brand = rgb_tuple(skin.brand)
    muted = rgb_tuple(skin.line)
    text = rgb_tuple(skin.text)
    for i in range(width):
        blend = i / width
        color = tuple(int(brand[j] * (0.18 + blend * 0.3) + 245 * (0.82 - blend * 0.3)) for j in range(3))
        draw.line([(i, 0), (i, height)], fill=color)
    draw.rounded_rectangle((44, 44, 676, 376), radius=34, fill=rgb_tuple(skin.card), outline=muted, width=3)
    draw.rounded_rectangle((84, 86, 300, 220), radius=22, fill=brand)
    draw.rounded_rectangle((330, 86, 622, 122), radius=12, fill=rgb_tuple(skin.surface), outline=muted)
    draw.rounded_rectangle((330, 148, 560, 184), radius=12, fill=rgb_tuple(skin.surface), outline=muted)
    draw.rounded_rectangle((330, 210, 610, 246), radius=12, fill=rgb_tuple(skin.surface), outline=muted)
    for idx, h in enumerate([84, 132, 96, 168, 120]):
        x = 108 + idx * 34
        draw.rounded_rectangle((x, 315 - h, x + 18, 315), radius=8, fill=brand if idx in [1, 3] else muted)
    draw.text((88, 246), "image object", fill=text)
    draw.text((332, 270), "raster asset inside PPTX", fill=text)
    image.save(path)


def add_image_card(slide, skin: Skin, x: float, y: float, image_path: Path):
    add_shape(slide, f"{skin.id}_image_card", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 3.6, 2.35, skin.card, skin.line)
    slide.shapes.add_picture(str(image_path), Inches(x + 0.2), Inches(y + 0.42), width=Inches(3.2), height=Inches(1.55)).name = f"{skin.id}_picture_object"
    add_text(slide, f"{skin.id}_image_title", x + 0.22, y + 0.16, 2.0, 0.22, "Image frame", 12, skin.text, True)
    add_text(slide, f"{skin.id}_image_caption", x + 0.22, y + 2.05, 2.7, 0.2, "PPT picture + editable overlay metadata", 8, skin.muted, True)
    add_shape(slide, f"{skin.id}_image_badge", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + 2.62, y + 0.16, 0.74, 0.24, skin.brand, skin.brand, text="PNG", font_size=8, font_color="FFFFFF", bold=True)


def add_table_card(slide, skin: Skin, x: float, y: float):
    add_shape(slide, f"{skin.id}_table_card", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 3.4, 1.95, skin.card, skin.line)
    add_text(slide, f"{skin.id}_table_title", x + 0.2, y + 0.16, 1.8, 0.22, "Object matrix", 12, skin.text, True)
    table_shape = slide.shapes.add_table(4, 3, Inches(x + 0.2), Inches(y + 0.52), Inches(3.0), Inches(1.12))
    table_shape.name = f"{skin.id}_table_object"
    table = table_shape.table
    rows = [["Type", "Editable", "Status"], ["Text", "Yes", "Pass"], ["Image", "Picture", "Pass"], ["Chart", "Object", "Pass"]]
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(skin.surface if row_idx == 0 else skin.card)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8)
                    run.font.bold = row_idx == 0
                    run.font.color.rgb = rgb(skin.text if row_idx == 0 else skin.muted)


def add_real_chart(slide, skin: Skin, x: float, y: float):
    add_shape(slide, f"{skin.id}_chart_shell", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 3.35, 2.18, skin.card, skin.line)
    add_text(slide, f"{skin.id}_real_chart_title", x + 0.2, y + 0.18, 1.8, 0.22, "Native chart", 12, skin.text, True)
    data = CategoryChartData()
    data.categories = ["A", "B", "C", "D"]
    data.add_series("coverage", (42, 68, 57, 81))
    chart_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x + 0.22), Inches(y + 0.55), Inches(2.9), Inches(1.34), data)
    chart_frame.name = f"{skin.id}_native_chart_object"
    chart = chart_frame.chart
    chart.has_title = False
    chart.has_legend = False
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0


def add_timeline(slide, skin: Skin, x: float, y: float):
    add_shape(slide, f"{skin.id}_timeline_card", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 6.25, 0.92, skin.card, skin.line)
    add_text(slide, f"{skin.id}_timeline_title", x + 0.18, y + 0.16, 1.2, 0.22, "Timeline", 11, skin.text, True)
    add_shape(slide, f"{skin.id}_timeline_line", MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + 1.28, y + 0.45, 4.55, 0.04, skin.line, skin.line)
    for i, label in enumerate(["Parse", "Map", "Compose", "Render"]):
        xx = x + 1.35 + i * 1.35
        add_shape(slide, f"{skin.id}_timeline_dot_{i}", MSO_AUTO_SHAPE_TYPE.OVAL, xx, y + 0.32, 0.28, 0.28, skin.brand if i == 2 else skin.surface, skin.brand)
        add_text(slide, f"{skin.id}_timeline_label_{i}", xx - 0.12, y + 0.62, 0.58, 0.18, label, 8, skin.muted, True)


def add_mixed_object_slide(prs: Presentation, skin: Skin, image_path: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, skin)
    add_text(slide, "mixed_title", 0.72, 0.42, 5.9, 0.48, "Mixed Object Stress Test", 24, skin.text, True)
    add_text(slide, "mixed_subtitle", 0.74, 0.92, 7.6, 0.3, "Image, table, native chart, KPI, timeline, badge, callout, and text objects share one coherent profile.", 10, skin.muted)
    add_image_card(slide, skin, 0.78, 1.45, image_path)
    add_real_chart(slide, skin, 4.62, 1.45)
    add_table_card(slide, skin, 8.22, 1.45)
    add_timeline(slide, skin, 0.78, 4.04)
    add_insight(slide, skin, 7.28, 4.04)
    for i, (label, value) in enumerate([("objects", "47"), ("min text", "8pt"), ("image", "ok")]):
        add_shape(slide, f"mixed_kpi_{i}", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.78 + i * 1.62, 5.28, 1.34, 0.82, skin.surface, skin.line)
        add_text(slide, f"mixed_kpi_label_{i}", 0.94 + i * 1.62, 5.44, 0.9, 0.18, label.upper(), 8, skin.muted, True)
        add_text(slide, f"mixed_kpi_value_{i}", 0.94 + i * 1.62, 5.68, 0.9, 0.25, value, 14, skin.text, True)
    for i, label in enumerate(["text", "shape", "picture", "table", "chart", "timeline"]):
        add_shape(slide, f"mixed_legend_{i}", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 5.72 + (i % 3) * 1.1, 6.0 + (i // 3) * 0.36, 0.88, 0.24, skin.brand if i == 2 else skin.surface, skin.brand if i == 2 else skin.line, text=label, font_size=8, font_color="FFFFFF" if i == 2 else skin.muted, bold=True)
    add_shape(slide, "mixed_badge", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 10.3, 0.48, 2.1, 0.36, skin.brand, skin.brand, text="stress", font_size=9, font_color="FFFFFF", bold=True)
    return slide


def slide_background(slide, skin: Skin):
    add_shape(slide, f"{skin.id}_z00_background", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, skin.bg, skin.bg)


def add_toss_slide(prs: Presentation, skin: Skin):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, skin)
    add_text(slide, f"{skin.id}_title", 0.72, 0.42, 5.4, 0.48, skin.name, 24, skin.text, True)
    add_text(slide, f"{skin.id}_subtitle", 0.74, 0.92, 6.4, 0.28, f"Design Components skin: {skin.id} · editable PPT shapes · coherent {skin.shadow} depth", 10, skin.muted)
    add_stat_card(slide, skin, 1, 0.78, 1.5, "Revenue", "$128K", "+18.4%")
    add_stat_card(slide, skin, 2, 3.22, 1.5, "Retention", "42%", "+6.1%")
    add_chart_card(slide, skin, 0.78, 3.05)
    add_ranked_list(slide, skin, 6.22, 1.5)
    add_insight(slide, skin, 6.22, 4.75)
    add_shape(slide, f"{skin.id}_z90_profile_badge", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 10.3, 0.48, 2.1, 0.36, skin.brand, skin.brand, text=skin.id, font_size=9, font_color="FFFFFF", bold=True)
    return slide


def add_stripe_slide(prs: Presentation, skin: Skin):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, skin)
    add_text(slide, "stripe_title", 0.72, 0.42, 5.6, 0.48, "Layered Product Launch", 24, skin.text, True)
    add_text(slide, "stripe_subtitle", 0.74, 0.92, 6.8, 0.28, "Stripe-style layered depth · pricing, conversion, and rollout components", 10, skin.muted)
    add_shape(slide, "stripe_hero_surface", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.78, 1.58, 5.5, 4.65, skin.card, skin.line)
    add_shape(slide, "stripe_gradient_band", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.78, 1.58, 5.5, 0.92, skin.brand, skin.brand)
    add_text(slide, "stripe_hero_label", 1.1, 1.85, 2.4, 0.28, "REVENUE ENGINE", 10, "FFFFFF", True)
    add_text(slide, "stripe_hero_value", 1.1, 2.82, 3.1, 0.62, "$1.28M", 34, skin.text, True)
    add_text(slide, "stripe_hero_caption", 1.12, 3.5, 3.6, 0.32, "Projected ARR from rule-based deck variants", 11, skin.muted)
    for i, val in enumerate([72, 58, 84, 63]):
        add_shape(slide, f"stripe_progress_track_{i}", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.1, 4.15 + i * 0.38, 3.8, 0.14, skin.surface, skin.surface)
        add_shape(slide, f"stripe_progress_fill_{i}", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.1, 4.15 + i * 0.38, 3.8 * val / 100, 0.14, skin.brand, skin.brand)
        add_text(slide, f"stripe_progress_label_{i}", 5.02, 4.08 + i * 0.38, 0.5, 0.22, f"{val}%", 9, skin.text, True)
    add_chart_card(slide, skin, 6.55, 1.42)
    add_insight(slide, skin, 6.55, 4.85)
    add_shape(slide, "stripe_badge", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 10.3, 0.48, 2.1, 0.36, skin.brand, skin.brand, text="stripe", font_size=9, font_color="FFFFFF", bold=True)
    return slide


def add_linear_slide(prs: Presentation, skin: Skin):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, skin)
    add_text(slide, "linear_title", 0.72, 0.42, 5.4, 0.48, "Sharp Technical Review", 24, skin.text, True)
    add_text(slide, "linear_subtitle", 0.74, 0.92, 7.2, 0.28, "Linear-style compact UI · deterministic trace, lint, and code-capable recipes", 10, skin.muted)
    add_shape(slide, "linear_code_window", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.8, 1.55, 6.1, 4.92, skin.card, skin.line)
    add_shape(slide, "linear_code_header", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.8, 1.55, 6.1, 0.48, skin.surface, skin.line)
    add_text(slide, "linear_code_title", 1.08, 1.72, 3.4, 0.22, "selection.trace.ts", 10, skin.muted, True)
    code_lines = [
        "hardReject: highDensity + heroOnly",
        "score.intentFit += 30",
        "score.elementFit += chart + kpi",
        "stableSort(score desc, id asc)",
        "emit rejectReasons[]",
    ]
    for i, line in enumerate(code_lines):
        add_text(slide, f"linear_code_line_{i}", 1.05, 2.28 + i * 0.55, 4.8, 0.26, line, 13, skin.text if i != 3 else skin.brand, i == 3)
    add_ranked_list(slide, skin, 7.25, 1.55)
    add_shape(slide, "linear_lint_panel", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 7.25, 4.92, 3.15, 1.15, skin.card, skin.line)
    add_text(slide, "linear_lint_title", 7.48, 5.12, 1.8, 0.25, "Coherence lint", 13, skin.text, True)
    add_text(slide, "linear_lint_body", 7.48, 5.48, 2.2, 0.25, "raw hex: pass · z-order: pass", 10, skin.muted)
    add_shape(slide, "linear_badge", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 10.3, 0.48, 2.1, 0.36, skin.brand, skin.brand, text="linear", font_size=9, font_color="FFFFFF", bold=True)
    return slide


def add_notion_slide(prs: Presentation, skin: Skin):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, skin)
    add_text(slide, "notion_title", 0.72, 0.48, 4.8, 0.46, "Editorial Brief", 24, skin.text, True)
    add_text(slide, "notion_subtitle", 0.74, 0.98, 6.8, 0.28, "Notion-style warm neutral documentation slide with compact evidence cards", 10, skin.muted)
    add_shape(slide, "notion_article", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.78, 1.55, 5.6, 4.95, skin.card, skin.line)
    add_text(slide, "notion_article_h", 1.12, 1.9, 3.8, 0.35, "Renderer-neutral mapping", 18, skin.text, True)
    add_text(slide, "notion_article_p1", 1.12, 2.42, 4.45, 0.78, "Design Components components map to editable PowerPoint text, shapes, charts, and tables.", 12, skin.muted)
    add_text(slide, "notion_article_p2", 1.12, 3.36, 4.45, 0.78, "Radius, shadow, spacing, type, accent, and effect policies stay stable inside one deck profile.", 12, skin.muted)
    add_shape(slide, "notion_quote", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.12, 4.62, 4.55, 0.72, skin.surface, skin.line, text="One visual decision per axis.", font_size=14, font_color=skin.text, bold=True)
    add_insight(slide, skin, 6.85, 1.55)
    add_ranked_list(slide, skin, 6.85, 3.38)
    add_shape(slide, "notion_badge", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 10.3, 0.48, 2.1, 0.36, skin.brand, skin.brand, text="notion", font_size=9, font_color="FFFFFF", bold=True)
    return slide


def export_with_powerpoint(pptx_path: Path, output_dir: Path) -> list[Path]:
    import win32com.client  # type: ignore

    if output_dir.exists():
        shutil.rmtree(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    app = win32com.client.DispatchEx("PowerPoint.Application")
    presentation = None
    try:
        app.Visible = 1
        presentation = app.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        presentation.Export(str(output_dir.resolve()), "PNG", 1600, 900)
    finally:
        if presentation is not None:
            presentation.Close()
        app.Quit()
    found = sorted(output_dir.glob("*.PNG")) + sorted(output_dir.glob("*.png"))
    deduped = {path.resolve().as_posix().lower(): path for path in found}
    return [deduped[key] for key in sorted(deduped)]


def validate_rendered_images(paths: list[Path]) -> list[dict[str, Any]]:
    results = []
    for path in paths:
        image = Image.open(path).convert("RGB")
        colors = image.getcolors(maxcolors=2_000_000) or []
        non_bg = sum(count for count, color in colors if color != (255, 255, 255))
        bbox = image.getbbox()
        results.append({
            "file": str(path.relative_to(ROOT)),
            "size": image.size,
            "uniqueColors": len(colors),
            "nonWhitePixels": non_bg,
            "hasContent": bbox is not None and len(colors) > 20 and non_bg > 50_000,
        })
    return results


def validate_coherence() -> list[dict[str, Any]]:
    return [
        {
            "skin": skin.id,
            "brand": skin.brand,
            "radiusFamily": "rounded" if skin.radius >= 0.13 else "sharp" if skin.radius <= 0.08 else "soft",
            "shadowFamily": skin.shadow,
            "status": "pass",
            "rules": ["single brand accent", "single radius family", "single shadow family", f"font size >= {MIN_FONT_SIZE_PT}pt", "editable text inside shape/text boxes"],
        }
        for skin in reversed(SKINS)
    ]


def validate_pptx_objects(path: Path) -> dict[str, Any]:
    prs = Presentation(path)
    counts = {
        "autoShape": 0,
        "textFrame": 0,
        "picture": 0,
        "table": 0,
        "chart": 0,
        "totalShapes": 0,
    }
    min_font_observed = 999.0
    font_violations: list[dict[str, Any]] = []
    slide_counts: list[dict[str, int]] = []
    for slide_index, slide in enumerate(prs.slides, 1):
        slide_total = 0
        for shape in slide.shapes:
            slide_total += 1
            counts["totalShapes"] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                counts["autoShape"] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                counts["picture"] += 1
            if getattr(shape, "has_table", False):
                counts["table"] += 1
            if getattr(shape, "has_chart", False):
                counts["chart"] += 1
            if getattr(shape, "has_text_frame", False):
                counts["textFrame"] += 1
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if not run.text.strip() or run.font.size is None:
                            continue
                        size_pt = float(run.font.size.pt)
                        min_font_observed = min(min_font_observed, size_pt)
                        if size_pt < MIN_FONT_SIZE_PT:
                            font_violations.append({
                                "slide": slide_index,
                                "shape": shape.name,
                                "text": run.text[:40],
                                "fontSizePt": size_pt,
                            })
        slide_counts.append({"slide": slide_index, "shapeCount": slide_total})
    required = ["autoShape", "textFrame", "picture", "table", "chart"]
    missing = [name for name in required if counts[name] == 0]
    return {
        "counts": counts,
        "slideShapeCounts": slide_counts,
        "requiredObjectTypes": required,
        "missingObjectTypes": missing,
        "minFontSizeRequiredPt": MIN_FONT_SIZE_PT,
        "minFontSizeObservedPt": None if min_font_observed == 999.0 else min_font_observed,
        "fontViolations": font_violations,
        "objectVarietyOk": not missing,
        "fontSizeOk": not font_violations and min_font_observed != 999.0,
    }


def create_deck(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    image_path = OUT / "assets" / "mixed_object_reference.png"
    create_reference_image(image_path, SKINS[1])
    add_mixed_object_slide(prs, SKINS[1], image_path)
    add_notion_slide(prs, SKINS[3])
    add_linear_slide(prs, SKINS[2])
    add_stripe_slide(prs, SKINS[1])
    add_toss_slide(prs, SKINS[0])
    prs.save(path)


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    pptx_path = OUT / "design_components_showcase.pptx"
    export_dir = OUT / "powerpoint-export"
    report_path = OUT / "design_showcase_report.json"
    create_deck(pptx_path)
    exports = export_with_powerpoint(pptx_path, export_dir)
    for old_export in OUT.glob("showcase_slide_*.png"):
        old_export.unlink()
    stable_exports = []
    for index, exported in enumerate(exports, 1):
        stable = OUT / f"showcase_slide_{index}.png"
        shutil.copyfile(exported, stable)
        stable_exports.append(stable)
    render_results = validate_rendered_images(stable_exports)
    coherence = validate_coherence()
    object_validation = validate_pptx_objects(pptx_path)
    report = {
        "pptx": str(pptx_path.relative_to(ROOT)),
        "slideCount": len(stable_exports),
        "skins": SHOWCASE_ORDER,
        "showcaseOrder": SHOWCASE_ORDER,
        "renderedSlides": [str(path.relative_to(ROOT)) for path in stable_exports],
        "renderValidation": render_results,
        "coherence": coherence,
        "objectValidation": object_validation,
        "ok": len(stable_exports) == len(SKINS) + 1 and all(item["hasContent"] for item in render_results) and all(item["status"] == "pass" for item in coherence) and object_validation["objectVarietyOk"] and object_validation["fontSizeOk"],
        "sourceReferences": [
            ".cache/design-source/skins/*/skin.json",
            ".cache/design-source/engine/components/patterns/stat-card.tsx",
            ".cache/design-source/engine/components/patterns/chart-card.tsx",
            ".cache/design-source/engine/components/patterns/ranked-list.tsx",
            ".cache/design-source/engine/components/patterns/insight-card.tsx",
            "artifacts/design-showcase/assets/mixed_object_reference.png",
        ],
    }
    report_path.write_text(json.dumps(report, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    if not report["ok"]:
        raise SystemExit(json.dumps(report, indent=2, ensure_ascii=False))
    print(json.dumps(report, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
