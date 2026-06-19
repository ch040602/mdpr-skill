#!/usr/bin/env python3
from __future__ import annotations

import json
import shutil
import subprocess
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
MDPR = ROOT / ".cache" / "mdpr"
SOURCE = ROOT / "artifacts" / "release-check" / "mdpr-skill-release-check.md"
OUT = ROOT / "artifacts" / "release-check"
BUILD = OUT / "build"
PPTX = OUT / "mdpr-skill-release-check.pptx"
PNG_DIR = OUT / "powerpoint-export"
REPORT = OUT / "mdpr-skill-release-check-report.json"
MIN_FONT_PT = 14.0


def build_deck() -> None:
    cli = MDPR / "packages" / "cli" / "dist" / "index.js"
    if BUILD.exists():
        shutil.rmtree(BUILD)
    BUILD.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        ["node", str(cli), "build", str(SOURCE), "--to", "pptx", "--out", str(BUILD), "--design", "editorial"],
        cwd=MDPR,
        check=True,
    )
    generated = BUILD / "deck.pptx"
    if not generated.is_file():
        raise FileNotFoundError(generated)
    shutil.copyfile(generated, PPTX)


def export_with_powerpoint() -> list[Path]:
    import win32com.client  # type: ignore

    if PNG_DIR.exists():
        shutil.rmtree(PNG_DIR)
    PNG_DIR.mkdir(parents=True, exist_ok=True)
    app = win32com.client.DispatchEx("PowerPoint.Application")
    presentation = None
    try:
        app.Visible = 1
        presentation = app.Presentations.Open(str(PPTX.resolve()), WithWindow=False)
        presentation.Export(str(PNG_DIR.resolve()), "PNG", 1600, 900)
    finally:
        if presentation is not None:
            presentation.Close()
        app.Quit()
    exported: dict[str, Path] = {}
    for path in list(PNG_DIR.glob("*.PNG")) + list(PNG_DIR.glob("*.png")):
        exported[str(path.resolve()).casefold()] = path
    return sorted(exported.values(), key=lambda item: item.name.casefold())


def validate_pptx() -> dict[str, Any]:
    prs = Presentation(PPTX)
    counts = {"slides": len(prs.slides), "shapes": 0, "textFrames": 0, "pictures": 0, "tables": 0, "charts": 0}
    min_font = 999.0
    for slide in prs.slides:
        for shape in slide.shapes:
            counts["shapes"] += 1
            if getattr(shape, "has_text_frame", False):
                counts["textFrames"] += 1
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip() and run.font.size is not None:
                            min_font = min(min_font, float(run.font.size.pt))
            if getattr(shape, "has_table", False):
                counts["tables"] += 1
            if getattr(shape, "has_chart", False):
                counts["charts"] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                counts["pictures"] += 1
    counts["minFontSizePt"] = None if min_font == 999.0 else min_font
    counts["fontFloorOk"] = counts["minFontSizePt"] is not None and counts["minFontSizePt"] >= MIN_FONT_PT
    return counts


def validate_pngs(paths: list[Path]) -> list[dict[str, Any]]:
    results = []
    for path in paths:
        image = Image.open(path).convert("RGB")
        colors = image.getcolors(maxcolors=4_000_000) or []
        non_white = sum(count for count, color in colors if color != (255, 255, 255))
        results.append(
            {
                "file": str(path.relative_to(ROOT)),
                "size": image.size,
                "uniqueColors": len(colors),
                "nonWhitePixels": non_white,
                "hasContent": len(colors) > 20 and non_white > 50_000,
            }
        )
    return results


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    build_deck()
    pngs = export_with_powerpoint()
    pptx_validation = validate_pptx()
    png_validation = validate_pngs(pngs)
    report = {
        "source": str(SOURCE.relative_to(ROOT)),
        "pptx": str(PPTX.relative_to(ROOT)),
        "pngDir": str(PNG_DIR.relative_to(ROOT)),
        "minRequiredFontPt": MIN_FONT_PT,
        "pptxValidation": pptx_validation,
        "pngValidation": png_validation,
        "ok": bool(pptx_validation["fontFloorOk"]) and bool(png_validation) and all(item["hasContent"] for item in png_validation),
    }
    REPORT.write_text(json.dumps(report, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    if not report["ok"]:
        raise SystemExit(json.dumps(report, indent=2, ensure_ascii=False))
    print(json.dumps(report, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
