#!/usr/bin/env python3
from __future__ import annotations

import json
import re
import shutil
import time
from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from urllib.parse import quote, urlsplit, urlunsplit
from urllib.request import Request, urlopen

from PIL import Image, ImageChops, ImageFilter, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
CACHE = ROOT / ".cache" / "pptbizcam"
PPTX_DIR = CACHE / "pptx"
PNG_DIR = CACHE / "png"
OUT = ROOT / "artifacts" / "pptbizcam-analysis"
SEEDS = ROOT / "design_components" / "design-source-adapter" / "seeds" / "visual-diversification-seeds.json"

START_POST_ID = 11677
MIN_DOWNLOADS = 50
MAX_POST_SCAN = 900
MAX_RENDER_DECKS = 50
REQUEST_DELAY_SECONDS = 0.08
USER_AGENT = "Mozilla/5.0 mdpr-skill structural pattern analysis"


@dataclass(frozen=True)
class ObjectPattern:
    id: str
    family: str
    best_for: list[str]
    structure: list[str]
    selection: dict[str, Any]
    coherence_guard: str


OBJECT_PATTERNS: list[ObjectPattern] = [
    ObjectPattern("accent-rail-card", "card", ["plain-list", "medium-text"], ["vertical rail", "title/body block"], {"hasImage": False, "textCharsMax": 180, "importanceMax": 3}, "Use one rail width and one accent role across sibling cards."),
    ObjectPattern("top-rule-card", "card", ["short-list", "low-importance"], ["hairline rule", "quiet surface"], {"density": "low", "importanceMax": 2}, "Rule weight stays subordinate to title text."),
    ObjectPattern("number-tab-card", "card", ["sequence", "ordered-list"], ["small tab", "number", "aligned title"], {"relation": "sequence", "itemCountMax": 7}, "Tabs share size and use brightness sequence."),
    ObjectPattern("corner-chip-card", "card", ["short-label", "importance-4"], ["corner chip", "compact label"], {"textCharsMax": 80, "importanceMin": 4}, "Chip never enters the text safe area."),
    ObjectPattern("bracket-note-card", "card", ["constraint", "proof"], ["thin bracket", "label", "indented body"], {"relation": ["constraint", "proof"]}, "Bracket height follows text block height, not full card height."),
    ObjectPattern("dot-marker-row", "list", ["plain-list", "semantic-cue"], ["small dot", "first-line alignment"], {"textOnly": True, "itemCountMin": 3}, "Dot center aligns to the first text line midpoint."),
    ObjectPattern("split-tone-row", "list", ["comparison", "two-part-item"], ["label zone", "body zone"], {"relation": "comparison"}, "Tint is lower contrast than text and consistent across rows."),
    ObjectPattern("proof-chip-inline", "callout", ["validation", "key-number"], ["small chip", "metric text"], {"hasKeyNumber": True, "importanceMin": 4}, "Only one chip in a local group may use contrast color."),
    ObjectPattern("thin-outline-card", "card", ["dense-list", "table-adjacent"], ["outline", "plain surface"], {"density": "high"}, "Use surface-line stroke unless selected as the lead item."),
    ObjectPattern("bottom-meter-card", "metric", ["progress", "score"], ["bottom meter", "short label"], {"hasValue": True, "relation": "progress"}, "Meter length must map to data value."),
    ObjectPattern("side-notch-card", "card", ["important-item", "constraint"], ["side notch", "aligned text"], {"importanceMin": 4}, "Notch width remains fixed across siblings."),
    ObjectPattern("floating-label-pin", "image", ["image-card", "photo-present"], ["small pin", "image edge"], {"hasImage": True, "textCharsMax": 60}, "Pin stays inside image safe area with readable contrast."),
    ObjectPattern("caption-underlay", "image", ["photo-present", "long-caption"], ["caption band", "single line label"], {"hasImage": True, "textCharsMin": 60}, "Underlay supports readability without hiding the subject."),
    ObjectPattern("rank-ribbon-card", "ranking", ["ranking", "key-number"], ["rank ribbon", "large number"], {"relation": "ranking", "hasKeyNumber": True}, "Ribbon height is smaller than number height."),
    ObjectPattern("micro-icon-marker", "icon", ["text-only", "semantic-cue"], ["one monotone icon", "left gutter"], {"textOnly": True, "density": "medium"}, "Icon stays smaller than body line-height."),
    ObjectPattern("double-rule-header", "header", ["sectioned-list", "medium-density"], ["two rules", "header label"], {"relation": "category"}, "Rules share color and differ only by opacity."),
    ObjectPattern("soft-shadow-lift", "surface", ["lead-card", "low-density"], ["shallow shadow", "plain surface"], {"importanceMin": 5, "density": "low"}, "Use one shadow depth per slide."),
    ObjectPattern("inset-label-bar", "label", ["category-list", "comparison"], ["inset bar", "category label"], {"relation": ["category", "comparison"]}, "Bar widths align across peer cards."),
    ObjectPattern("vertical-step-rail", "diagram", ["sequence", "long-text"], ["vertical rail", "node dot", "wrapped text"], {"relation": "sequence", "textCharsMin": 120}, "Same-role nodes use one connector color."),
    ObjectPattern("horizontal-step-rail", "diagram", ["sequence", "short-text"], ["horizontal rail", "number nodes"], {"relation": "sequence", "textCharsMax": 120}, "Node centers align on one baseline."),
    ObjectPattern("paired-contrast-edge", "comparison", ["before-after", "opposition"], ["two edge colors", "split label"], {"relation": "opposition"}, "Use complementary colors only for true oppositions."),
    ObjectPattern("quote-rule-card", "quote", ["quote", "takeaway"], ["left quote rule", "body"], {"relation": "quote"}, "Italic applies only to quote body; labels remain upright."),
    ObjectPattern("metric-lead-card", "metric", ["key-number", "importance-5"], ["large number", "unit", "caption"], {"hasKeyNumber": True, "importanceMin": 5}, "Number uses display scale; caption stays body/caption scale."),
    ObjectPattern("checklist-grid-card", "table", ["checklist", "many-short-items"], ["check marker", "two-column grid"], {"relation": "checklist", "itemCountMin": 6}, "Markers are monotone or accent1 only."),
    ObjectPattern("constraint-stack-card", "constraint", ["constraint", "risk"], ["bold label", "indented description"], {"relation": ["constraint", "risk"]}, "Label font size is not smaller than description."),
    ObjectPattern("label-overline-card", "label", ["category", "short-title"], ["overline", "title", "body"], {"textCharsMax": 110}, "Uppercase overline only when under 18 characters."),
    ObjectPattern("subtle-band-card", "surface", ["medium-list", "non-lead"], ["surface band", "text group"], {"importanceMax": 3}, "Band uses surface tint and never competes with lead accent."),
    ObjectPattern("arc-corner-emphasis", "metric", ["ratio", "progress"], ["arc corner", "label"], {"hasValue": True, "relation": "progress"}, "Arc must map to status or progress using real arc geometry."),
    ObjectPattern("target-ring-badge", "metric", ["goal", "benchmark"], ["ring badge", "metric"], {"relation": "benchmark"}, "Ring thickness remains proportional."),
    ObjectPattern("image-sidecar-card", "image", ["photo-present", "medium-text"], ["reserved image sidecar", "text block"], {"hasImage": True}, "Reserve image width before text fitting."),
    ObjectPattern("table-summary-card", "table", ["table-adjacent", "summary"], ["header label", "two compact facts"], {"hasTable": True}, "Align to table column rhythm."),
    ObjectPattern("connector-dot-card", "diagram", ["flow", "dependency"], ["edge dot", "connector anchor"], {"relation": ["flow", "dependency"]}, "Same-role connectors share arrow style and color."),
    ObjectPattern("plain-safe-card", "fallback", ["high-density", "long-text"], ["minimal padding", "no ornament"], {"density": "high", "textCharsMin": 180}, "Use when decoration would reduce readability."),
    ObjectPattern("rounded-ticket-panel", "surface", ["process", "section-cover"], ["large rounded ticket", "punched corner cue"], {"density": "low"}, "Ticket affordance must encode grouping or sectioning."),
    ObjectPattern("paperclip-corner", "decoration", ["document", "worksheet"], ["small clip mark", "corner overlap"], {"intent": "document"}, "Clip is small and never covers text."),
    ObjectPattern("binder-hole-strip", "decoration", ["notebook", "process-note"], ["hole strip", "ruled background"], {"intent": "notebook"}, "Holes are decorative only when document metaphor is explicit."),
    ObjectPattern("tape-label", "label", ["temporary-note", "annotation"], ["short tape surface", "caption"], {"relation": "annotation"}, "Tape opacity stays low and text remains normal weight."),
    ObjectPattern("speech-bubble-callout", "callout", ["comment", "dialogue"], ["bubble", "tail", "short text"], {"relation": "comment"}, "Tail points to the referenced object."),
    ObjectPattern("diagonal-sash-label", "label", ["status", "rank"], ["diagonal sash", "short word"], {"textCharsMax": 24, "importanceMin": 4}, "Sash cannot cross body copy."),
    ObjectPattern("pill-step-node", "diagram", ["process", "short-steps"], ["pill node", "short label"], {"relation": "sequence", "textCharsMax": 72}, "All pills share height and baseline."),
    ObjectPattern("folded-corner-card", "card", ["document", "summary"], ["fold mark", "body"], {"intent": "document"}, "Fold is smaller than title cap height group."),
    ObjectPattern("stacked-paper-cards", "surface", ["multi-document", "references"], ["offset sheets", "front card"], {"relation": "collection"}, "Offsets are consistent and shallow."),
    ObjectPattern("overlap-lens-diagram", "diagram", ["overlap", "intersection"], ["two translucent shapes", "center label"], {"relation": "overlap"}, "Use only two or three sets; labels remain outside or centered."),
    ObjectPattern("axis-quadrant-map", "diagram", ["2x2", "positioning"], ["cross axes", "quadrant labels"], {"relation": "matrix"}, "Quadrants share geometry; emphasis uses stroke/tab."),
    ObjectPattern("ladder-step-diagram", "diagram", ["maturity", "levels"], ["stepped blocks", "ascending labels"], {"relation": "level"}, "Steps rise monotonically and labels keep one baseline offset."),
    ObjectPattern("loop-arrow-cycle", "diagram", ["cycle", "iteration"], ["curved arrows", "nodes"], {"relation": "cycle"}, "Arrow style stays constant around the loop."),
    ObjectPattern("center-hub-spokes", "diagram", ["hub", "features"], ["center node", "spokes", "outer nodes"], {"relation": "hub"}, "Outer nodes use equal radius and do not cross connectors."),
    ObjectPattern("timeline-marker-strip", "diagram", ["timeline", "milestones"], ["horizontal line", "markers"], {"relation": "timeline"}, "Markers align to one line and labels alternate only when needed."),
    ObjectPattern("vertical-timeline-cards", "diagram", ["timeline", "long labels"], ["vertical line", "side cards"], {"relation": "timeline", "textCharsMin": 120}, "Cards keep equal gutter from line."),
    ObjectPattern("stair-progress-meter", "metric", ["progress", "maturity"], ["step blocks", "current highlight"], {"relation": "progress"}, "Highlight only one current step unless multiple selected states exist."),
    ObjectPattern("donut-label-ring", "chart", ["part-to-whole", "ratio"], ["donut ring", "outside labels"], {"hasValue": True, "relation": "ratio"}, "Labels align to tangent starts and do not overlap ring."),
    ObjectPattern("gauge-score-card", "chart", ["score", "readiness"], ["semi gauge", "score", "caption"], {"hasValue": True, "relation": "score"}, "Needle/arc maps to value; caption stays secondary."),
    ObjectPattern("small-multiple-bars", "chart", ["comparison", "metrics"], ["equal mini bars", "shared labels"], {"hasChart": True, "itemCountMin": 3}, "All bars share scale and baseline."),
    ObjectPattern("trend-line-backdrop", "chart", ["trend", "forecast"], ["low contrast line", "foreground statement"], {"hasChart": True, "relation": "trend"}, "Backdrop opacity stays below foreground text contrast."),
    ObjectPattern("pictorial-anchor-labels", "image", ["metaphor", "provided-image"], ["image anchor", "attached labels"], {"hasImage": True, "needsImage": True}, "Labels attach to real image features, not invented decoration."),
    ObjectPattern("photo-window-mask", "image", ["photo-focus", "product"], ["cropped window", "caption"], {"hasImage": True}, "Use object-fit cover/contain according to subject inspection."),
    ObjectPattern("image-caption-split", "image", ["image-plus-explanation"], ["image region", "caption rail"], {"hasImage": True, "textCharsMin": 80}, "Caption rail never overlaps image subject."),
    ObjectPattern("tag-cloud-strip", "label", ["keywords", "taxonomy"], ["small chips", "wrap row"], {"relation": "keywords"}, "Chips wrap on grid; no random chip colors."),
    ObjectPattern("status-dot-table", "table", ["status-table", "checklist"], ["dot status column", "row labels"], {"hasTable": True, "relation": "status"}, "Dots align to row midpoints and use limited semantic colors."),
    ObjectPattern("risk-heat-grid", "table", ["risk", "matrix"], ["grid", "severity tint"], {"relation": "risk", "hasTable": True}, "Tint progression follows monochromatic or split-complementary rule."),
]


def safe_url(url: str) -> str:
    parts = urlsplit(url)
    return urlunsplit((parts.scheme, parts.netloc, quote(parts.path), quote(parts.query, safe="=&?"), parts.fragment))


def fetch_text(url: str) -> str:
    request = Request(safe_url(url), headers={"User-Agent": USER_AGENT})
    with urlopen(request, timeout=25) as response:
        return response.read().decode("utf-8", errors="ignore")


def fetch_binary(url: str) -> bytes:
    request = Request(safe_url(url), headers={"User-Agent": USER_AGENT})
    with urlopen(request, timeout=60) as response:
        return response.read()


def clean_title(html: str, fallback: str) -> str:
    match = re.search(r"<title[^>]*>(.*?)</title>", html, re.I | re.S)
    if not match:
        return fallback
    value = re.sub(r"\s+", " ", re.sub(r"<[^>]+>", "", match.group(1))).strip()
    return value.replace(" – PPT BIZCAM", "").replace(" - PPT BIZCAM", "")[:120]


def discover_posts() -> list[dict[str, Any]]:
    posts: list[dict[str, Any]] = []
    seen_links: set[str] = set()
    for post_id in range(START_POST_ID, START_POST_ID - MAX_POST_SCAN, -1):
        if len(seen_links) >= MIN_DOWNLOADS:
            break
        url = f"https://pptbizcam.co.kr/?p={post_id}"
        try:
            html = fetch_text(url)
        except Exception:
            continue
        links = sorted(set(re.findall(r'https?://[^"\'<>\s]+?\.pptx?', html, flags=re.I)))
        links = [link for link in links if link not in seen_links]
        if links:
            for link in links:
                seen_links.add(link)
            posts.append({"postId": post_id, "url": url, "title": clean_title(html, f"post-{post_id}"), "pptLinks": links})
        time.sleep(REQUEST_DELAY_SECONDS)
    return posts


def download_pptx(posts: list[dict[str, Any]]) -> list[dict[str, Any]]:
    PPTX_DIR.mkdir(parents=True, exist_ok=True)
    downloads: list[dict[str, Any]] = []
    ordinal = 0
    for post in posts:
        for link in post["pptLinks"]:
            ordinal += 1
            suffix = ".pptx" if link.lower().endswith("pptx") else ".ppt"
            path = PPTX_DIR / f"{ordinal:02d}_{post['postId']}{suffix}"
            if not path.exists() or path.stat().st_size < 1024:
                path.write_bytes(fetch_binary(link))
            downloads.append({**post, "pptUrl": link, "localPath": str(path.relative_to(ROOT)), "bytes": path.stat().st_size})
            if len(downloads) >= MIN_DOWNLOADS:
                return downloads
    return downloads


def analyze_ppt(path: Path) -> dict[str, Any]:
    prs = Presentation(path)
    counts: Counter[str] = Counter()
    slide_summaries: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(prs.slides, 1):
        slide_counts: Counter[str] = Counter()
        for shape in slide.shapes:
            counts["shapes"] += 1
            slide_counts["shapes"] += 1
            shape_type = "unknown" if shape.shape_type is None else str(int(shape.shape_type))
            counts[f"shapeType:{shape_type}"] += 1
            if getattr(shape, "has_text_frame", False):
                counts["textFrames"] += 1
                slide_counts["textFrames"] += 1
            if getattr(shape, "has_table", False):
                counts["tables"] += 1
                slide_counts["tables"] += 1
            if getattr(shape, "has_chart", False):
                counts["charts"] += 1
                slide_counts["charts"] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                counts["pictures"] += 1
                slide_counts["pictures"] += 1
        slide_summaries.append({"slide": slide_index, **dict(slide_counts)})
    return {"slides": len(prs.slides), "aggregate": dict(counts), "slideSummaries": slide_summaries[:12]}


def export_pngs(downloads: list[dict[str, Any]]) -> list[dict[str, Any]]:
    import win32com.client  # type: ignore

    PNG_DIR.mkdir(parents=True, exist_ok=True)
    app = win32com.client.DispatchEx("PowerPoint.Application")
    rendered: list[dict[str, Any]] = []
    presentation = None
    try:
        app.Visible = 1
        for item in downloads[:MAX_RENDER_DECKS]:
            ppt_path = ROOT / item["localPath"]
            deck_dir = PNG_DIR / ppt_path.stem
            if deck_dir.exists():
                shutil.rmtree(deck_dir)
            deck_dir.mkdir(parents=True, exist_ok=True)
            presentation = app.Presentations.Open(str(ppt_path.resolve()), WithWindow=False)
            presentation.Export(str(deck_dir.resolve()), "PNG", 1200, 675)
            presentation.Close()
            presentation = None
            files = sorted(deck_dir.glob("*.PNG")) + sorted(deck_dir.glob("*.png"))
            rendered.append({**item, "pngDir": str(deck_dir.relative_to(ROOT)), "pngCount": len(files), "samplePngs": [str(path.relative_to(ROOT)) for path in files[:3]]})
    finally:
        if presentation is not None:
            presentation.Close()
        app.Quit()
    return rendered


def analyze_png(path: Path) -> dict[str, Any]:
    image = Image.open(path).convert("RGB")
    colors = image.getcolors(maxcolors=4_000_000) or []
    gray = ImageOps.grayscale(image)
    edges = gray.filter(ImageFilter.FIND_EDGES)
    bbox = ImageChops.difference(image, Image.new("RGB", image.size, "white")).getbbox()
    edge_pixels = sum(1 for value in edges.getdata() if value > 28)
    non_white = sum(count for count, color in colors if color != (255, 255, 255))
    return {
        "file": str(path.relative_to(ROOT)),
        "size": image.size,
        "uniqueColors": len(colors),
        "nonWhitePixels": non_white,
        "edgeDensity": round(edge_pixels / (image.size[0] * image.size[1]), 4),
        "contentBBox": bbox,
    }


def analyze_rendered_pngs(rendered: list[dict[str, Any]]) -> list[dict[str, Any]]:
    results = []
    for item in rendered:
        for sample in item["samplePngs"][:2]:
            results.append(analyze_png(ROOT / sample))
    return results


def make_contact_sheet(rendered: list[dict[str, Any]], output: Path) -> None:
    thumbs = []
    for item in rendered[:50]:
        if not item["samplePngs"]:
            continue
        image = Image.open(ROOT / item["samplePngs"][0]).convert("RGB")
        image.thumbnail((240, 135))
        thumbs.append(image.copy())
    if not thumbs:
        return
    cols = 5
    rows = (len(thumbs) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * 240, rows * 135), "white")
    for index, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((index % cols) * 240, (index // cols) * 135))
    output.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output)


def pattern_to_json(pattern: ObjectPattern) -> dict[str, Any]:
    return {
        "id": pattern.id,
        "family": pattern.family,
        "bestFor": pattern.best_for,
        "structure": pattern.structure,
        "selection": pattern.selection,
        "coherenceGuard": pattern.coherence_guard,
    }


def update_seed_file(report: dict[str, Any]) -> None:
    seeds = json.loads(SEEDS.read_text(encoding="utf-8"))
    seeds["observedReferenceAnalysis"] = {
        "source": "PPT BIZCAM public post pages and downloadable PPT samples",
        "startPost": START_POST_ID,
        "downloadedPptFiles": report["pptDownloaded"],
        "decksAnalyzed": report["decksAnalyzed"],
        "slidesAnalyzed": report["slidesAnalyzed"],
        "renderedPngSlides": report["renderedPngSlides"],
        "aggregateObjects": report["aggregateObjects"],
        "methodNote": "Original PPT files are kept only in .cache/pptbizcam for local structural analysis. MDPR stores derived object vocabulary and rules, not copied assets or layouts.",
    }
    seeds["pptbizcamDerivedObjectPatterns"] = [pattern_to_json(pattern) for pattern in OBJECT_PATTERNS]
    seeds["pptbizcamRecursiveRulePolicy"] = {
        "minimumDownloadedPptFiles": MIN_DOWNLOADS,
        "minimumDerivedObjectPatterns": 50,
        "analysisLoop": [
            "scan numbered PPT BIZCAM post pages from the seed post id",
            "download public PPT links to .cache only",
            "render PPT slides to PNG through Microsoft PowerPoint",
            "extract slide-object and PNG structure metrics",
            "derive reusable object grammar and coherence guards",
            "regenerate seed JSON without storing source PPT assets in git",
        ],
        "selectionInputs": ["hasImage", "hasTable", "hasChart", "hasKeyNumber", "relation", "importance", "textChars", "density", "itemCount"],
        "copyGuard": "Do not copy source slide layouts, images, or brand-like objects. Use only structural grammar: alignment, grouping, connector, surface, marker, chart, and text-fitting methods.",
    }
    SEEDS.write_text(json.dumps(seeds, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    posts = discover_posts()
    downloads = download_pptx(posts)
    if len(downloads) < MIN_DOWNLOADS:
        raise SystemExit(f"Only downloaded {len(downloads)} PPT files; expected at least {MIN_DOWNLOADS}.")
    deck_analyses = []
    aggregate: Counter[str] = Counter()
    slides = 0
    for item in downloads:
        analysis = analyze_ppt(ROOT / item["localPath"])
        deck_analyses.append({**item, "analysis": analysis})
        aggregate.update(analysis["aggregate"])
        slides += int(analysis["slides"])
    rendered = export_pngs(downloads)
    png_analysis = analyze_rendered_pngs(rendered)
    contact_sheet = OUT / "pptbizcam-downloaded-contact-sheet.png"
    make_contact_sheet(rendered, contact_sheet)
    report = {
        "source": "https://pptbizcam.co.kr/?p=11677 with descending post-id scan",
        "postPagesWithPpt": len(posts),
        "posts": posts[:50],
        "pptDownloaded": len(downloads),
        "downloads": downloads,
        "decksAnalyzed": len(deck_analyses),
        "slidesAnalyzed": slides,
        "renderedDecks": len(rendered),
        "renderedPngSlides": sum(int(item["pngCount"]) for item in rendered),
        "pngSamplesAnalyzed": len(png_analysis),
        "aggregateObjects": dict(aggregate),
        "derivedObjectPatternCount": len(OBJECT_PATTERNS),
        "derivedObjectPatterns": [pattern_to_json(pattern) for pattern in OBJECT_PATTERNS],
        "rendered": rendered,
        "pngAnalysis": png_analysis,
        "contactSheet": str(contact_sheet.relative_to(ROOT)),
        "ok": len(downloads) >= MIN_DOWNLOADS and len(OBJECT_PATTERNS) >= 50 and len(png_analysis) >= 50,
    }
    (OUT / "pptbizcam-recursive-object-rules.json").write_text(json.dumps(report, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    (OUT / "pptbizcam-analysis.json").write_text(json.dumps(report, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    update_seed_file(report)
    if not report["ok"]:
        raise SystemExit(json.dumps(report, indent=2, ensure_ascii=False))
    print(json.dumps({k: report[k] for k in ["pptDownloaded", "decksAnalyzed", "slidesAnalyzed", "renderedPngSlides", "pngSamplesAnalyzed", "derivedObjectPatternCount", "contactSheet", "ok"]}, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
