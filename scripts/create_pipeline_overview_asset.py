#!/usr/bin/env python3
from __future__ import annotations

import json
import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "assets"
PPTX = OUT / "pipeline-overview.pptx"
PNG = OUT / "pipeline-overview.png"
REPORT = OUT / "pipeline-overview-report.json"


@dataclass(frozen=True)
class PipelineLayout:
    slide_w: float = 13.333
    slide_h: float = 7.25
    export_w: int = 2600
    export_h: int = 1414
    origin_x: float = 0.64
    origin_y: float = 0.34
    canvas_x: float = 0.54
    canvas_y: float = 1.0
    canvas_w: float = 12.24
    canvas_h: float = 5.78
    scale: float = 1.0
    min_font_pt: int = 8


LAYOUT = PipelineLayout()
TEXT_BOUNDS: list[dict[str, float | str]] = []
BOX_BOUNDS: dict[str, tuple[float, float, float, float]] = {}


def rgb(hex_value: str) -> RGBColor:
    return RGBColor(*(int(hex_value[i:i + 2], 16) for i in (0, 2, 4)))


def font_size(size: int) -> int:
    return max(LAYOUT.min_font_pt, round(size * LAYOUT.scale))


def fits_inside(inner: tuple[float, float, float, float], outer: tuple[float, float, float, float]) -> bool:
    ix, iy, iw, ih = inner
    ox, oy, ow, oh = outer
    eps = 0.012
    return ix + eps >= ox and iy + eps >= oy and ix + iw <= ox + ow + eps and iy + ih <= oy + oh + eps


def shape(slide, name: str, kind: MSO_AUTO_SHAPE_TYPE, x: float, y: float, w: float, h: float, fill: str, line: str = "D7E1EE"):
    item = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    item.name = name
    item.fill.solid()
    item.fill.fore_color.rgb = rgb(fill)
    item.line.color.rgb = rgb(line)
    BOX_BOUNDS[name] = (x, y, w, h)
    return item


def text(slide, name: str, x: float, y: float, w: float, h: float, value: str, size: int, color: str = "0F172A", bold: bool = False, parent: str | None = None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = value
    run.font.size = Pt(font_size(size))
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    TEXT_BOUNDS.append({"name": name, "parent": parent or "", "x": x, "y": y, "w": w, "h": h, "fontSize": font_size(size), "text": value})
    return box


def card(slide, name: str, x: float, y: float, w: float, h: float, title: str, lines: list[str], accent: str = "E2E8F0", border: str = "CBD5E1"):
    parent = f"{name}_card"
    shape(slide, parent, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, "F8FAFC", border)
    shape(slide, f"{name}_dot", MSO_AUTO_SHAPE_TYPE.OVAL, x + 0.16, y + 0.18, 0.18, 0.18, accent, accent)
    text(slide, f"{name}_title", x + 0.42, y + 0.13, w - 0.58, 0.22, title, 11, "111827", True, parent=parent)
    body_top = y + 0.44
    body_bottom = y + h - 0.12
    line_count = max(1, len(lines))
    line_gap = min(0.18, max(0.12, (body_bottom - body_top) / line_count))
    line_h = min(0.16, max(0.12, line_gap - 0.02))
    for idx, line in enumerate(lines):
        line_y = body_top + idx * line_gap
        text(slide, f"{name}_line_{idx}", x + 0.22, line_y, w - 0.42, line_h, line, 8, "526071", parent=parent)


def panel(slide, name: str, x: float, y: float, w: float, h: float, title: str, subtitle: str, fill: str, border: str, title_color: str):
    parent = f"{name}_panel"
    shape(slide, parent, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill, border)
    text(slide, f"{name}_title", x + 0.18, y + 0.18, w - 0.36, 0.22, title, 10, title_color, True, parent=parent)
    text(slide, f"{name}_subtitle", x + 0.18, y + 0.43, w - 0.36, 0.18, subtitle, 8, "64748B", parent=parent)
    return parent


def arrow_between(slide, name: str, x1: float, y1: float, x2: float, y2: float, color: str = "475569", width: float = 1.8, dashed: bool = False):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.name = name
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    if dashed:
        line.line.dash_style = 4
    return line


def arrow(slide, name: str, x: float, y: float, w: float, color: str = "475569", dashed: bool = False):
    return arrow_between(slide, name, x, y, x + w, y, color, 1.8, dashed)


def build_deck(path: Path) -> None:
    TEXT_BOUNDS.clear()
    BOX_BOUNDS.clear()
    prs = Presentation()
    prs.slide_width = Inches(LAYOUT.slide_w)
    prs.slide_height = Inches(LAYOUT.slide_h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    shape(slide, "z00_background", MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, LAYOUT.slide_w, LAYOUT.slide_h, "F8FAFC", "F8FAFC")
    shape(slide, "z01_canvas", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, LAYOUT.canvas_x, LAYOUT.canvas_y, LAYOUT.canvas_w, LAYOUT.canvas_h, "FFFFFF", "DDE6F3")

    text(slide, "title", LAYOUT.origin_x, LAYOUT.origin_y, 6.7, 0.42, "MDPR Design Components Pipeline", 23, "0F172A", True)
    text(slide, "subtitle", LAYOUT.origin_x, 0.7, 8.85, 0.24, "Content is split by MDPR. LLM reasoning supplies hints. Deterministic rules own layout, styling, z-order, and editable rendering.", 9, "475569")

    panel(slide, "zone_content", 0.82, 1.22, 2.55, 3.82, "1. Content Contract", "MDPR creates semantic structure only", "F8FAFC", "D7E1EE", "334155")
    panel(slide, "zone_reasoning", 3.58, 1.22, 2.56, 3.82, "2. LLM Reasoning", "optional intent and grouping hints", "EFF6FF", "93C5FD", "1D4ED8")
    panel(slide, "zone_rules", 6.36, 1.22, 3.62, 3.82, "3. Deterministic Design", "final visual choices happen here", "EEF7F1", "86EFAC", "15803D")
    panel(slide, "zone_outputs", 10.2, 1.22, 2.24, 3.82, "4. Editable Outputs", "PPTX, HTML, PDF, and checks", "F8FAFC", "D7E1EE", "334155")

    shape(slide, "start_flag", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.05, 1.78, 0.74, 0.26, "111827", "111827")
    text(slide, "start_flag_text", 1.2, 1.86, 0.44, 0.1, "start", 8, "FFFFFF", True, parent="start_flag")
    card(slide, "markdown", 1.02, 2.16, 1.98, 0.92, "Markdown", ["text, tables, code", "images and notes"], "CBD5E1")
    card(slide, "splitter", 1.02, 3.42, 1.98, 1.02, "MDPR Splitter", ["slide and object split", "no visual choices"], "CBD5E1")
    arrow_between(slide, "main_start_markdown", 1.4, 2.04, 1.4, 2.16, "111827", 2.4)
    arrow_between(slide, "main_markdown_splitter", 2.02, 3.08, 2.02, 3.42, "475569", 3.0)

    shape(slide, "ir_core", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 3.86, 1.78, 1.96, 0.94, "FFFFFF", "93C5FD")
    text(slide, "ir_core_title", 4.08, 1.98, 1.34, 0.2, "Slide Element IR", 12, "111827", True, parent="ir_core")
    text(slide, "ir_core_body", 4.08, 2.28, 1.44, 0.16, "content-only object contract", 8, "526071", parent="ir_core")
    shape(slide, "reasoning_card", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 3.84, 3.0, 2.0, 1.42, "FFFFFF", "93C5FD")
    text(slide, "reasoning_title", 4.06, 3.18, 1.42, 0.24, "Reasoning Result", 12, "111827", True, parent="reasoning_card")
    text(slide, "reasoning_body", 4.06, 3.52, 1.46, 0.2, "intent, grouping, importance", 8, "526071", parent="reasoning_card")
    shape(slide, "reasoning_badge", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 4.06, 3.88, 1.48, 0.24, "DBEAFE", "93C5FD")
    text(slide, "reasoning_badge_text", 4.24, 3.96, 1.08, 0.1, "hints only", 8, "1D4ED8", True, parent="reasoning_badge")
    text(slide, "reasoning_guardrail", 4.06, 4.2, 1.48, 0.16, "no coordinates or styles", 8, "64748B", True, parent="reasoning_card")

    arrow_between(slide, "main_splitter_ir", 3.0, 3.92, 3.86, 2.26, "475569", 3.4)
    arrow_between(slide, "hint_ir_reasoning", 4.86, 2.72, 4.86, 3.0, "2563EB", 2.0, True)

    shape(slide, "rule_engine", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 6.66, 1.76, 2.92, 0.74, "DCFCE7", "86EFAC")
    text(slide, "rule_engine_title", 6.9, 1.94, 2.0, 0.22, "Rule Engine Boundary", 13, "14532D", True, parent="rule_engine")
    text(slide, "rule_engine_body", 6.9, 2.25, 2.22, 0.12, "recipes, variants, coordinates, z-order", 8, "166534", parent="rule_engine")
    card(slide, "features", 6.66, 2.84, 1.36, 0.78, "Features", ["density, mix", "size risk"], "86EFAC", "BBF7D0")
    card(slide, "recipes", 8.28, 2.84, 1.36, 0.78, "Recipes", ["profile match", "component variant"], "86EFAC", "BBF7D0")
    card(slide, "compose", 6.66, 3.9, 1.36, 0.78, "Compose", ["regions, fit", "overflow policy"], "86EFAC", "BBF7D0")
    card(slide, "decorate", 8.28, 3.9, 1.36, 0.78, "Decorate", ["type, radius", "shadow, effects"], "86EFAC", "BBF7D0")
    arrow_between(slide, "main_ir_rules", 5.82, 2.18, 6.66, 2.18, "111827", 4.2)
    arrow_between(slide, "hint_reasoning_rules", 5.84, 3.66, 6.66, 3.24, "2563EB", 2.2, True)
    arrow_between(slide, "rule_features_recipes", 8.02, 3.23, 8.28, 3.23, "16A34A", 1.7)
    arrow_between(slide, "rule_features_compose", 7.34, 3.62, 7.34, 3.9, "16A34A", 1.4)
    arrow_between(slide, "rule_recipes_decorate", 8.96, 3.62, 8.96, 3.9, "16A34A", 1.4)
    arrow_between(slide, "rule_compose_decorate", 8.02, 4.29, 8.28, 4.29, "16A34A", 1.7)

    card(slide, "styled_ir", 10.42, 1.82, 1.66, 0.92, "Styled Deck IR", ["renderer-neutral", "visual contract"], "CBD5E1")
    card(slide, "renderers", 10.42, 3.08, 1.66, 0.94, "Renderers", ["editable PPTX", "HTML and PDF"], "CBD5E1")
    shape(slide, "visual_check", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 10.42, 4.34, 1.66, 0.42, "FEF3C7", "F59E0B")
    text(slide, "visual_check_title", 10.62, 4.49, 1.18, 0.12, "visual validation", 8, "92400E", True, parent="visual_check")
    arrow_between(slide, "main_rules_styled_ir", 9.58, 2.18, 10.42, 2.18, "111827", 4.2)
    arrow_between(slide, "main_styled_renderers", 11.25, 2.74, 11.25, 3.08, "475569", 2.4)
    arrow_between(slide, "validation_loop", 11.25, 4.02, 11.25, 4.34, "F59E0B", 2.0)

    shape(slide, "coherence_band", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.88, 5.56, 11.54, 0.76, "F8FAFC", "D7E1EE")
    text(slide, "coherence_title", 1.14, 5.78, 1.68, 0.18, "Coherence checks", 12, "111827", True, parent="coherence_band")
    text(slide, "coherence_body", 2.88, 5.78, 6.9, 0.22, "One visual language across mixed objects: hierarchy-scaled type, bounded text, consistent spacing, aligned starts, and readable minimum sizes.", 8, "526071", parent="coherence_band")
    shape(slide, "font_badge", MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 10.48, 5.76, 1.42, 0.28, "111827", "111827")
    text(slide, "font_badge_text", 10.66, 5.85, 1.08, 0.1, "font scale by role", 8, "FFFFFF", True, parent="font_badge")

    prs.save(path)


def export_with_powerpoint(pptx_path: Path, output_dir: Path) -> Path:
    import win32com.client  # type: ignore

    if output_dir.exists():
        shutil.rmtree(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    app = win32com.client.DispatchEx("PowerPoint.Application")
    presentation = None
    try:
        app.Visible = 1
        presentation = app.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        presentation.Export(str(output_dir.resolve()), "PNG", LAYOUT.export_w, LAYOUT.export_h)
    finally:
        if presentation is not None:
            presentation.Close()
        app.Quit()
    candidates = sorted(output_dir.glob("*.PNG")) + sorted(output_dir.glob("*.png"))
    if not candidates:
        raise FileNotFoundError(f"PowerPoint did not export PNG files to {output_dir}")
    return candidates[0]


def validate_png(path: Path) -> dict[str, Any]:
    image = Image.open(path).convert("RGB")
    colors = image.getcolors(maxcolors=5_000_000) or []
    non_white = sum(count for count, color in colors if color != (255, 255, 255))
    return {
        "file": str(path.relative_to(ROOT)),
        "size": image.size,
        "uniqueColors": len(colors),
        "nonWhitePixels": non_white,
        "hasContent": len(colors) > 50 and non_white > 100_000,
    }


def validate_layout() -> dict[str, Any]:
    overflow: list[dict[str, Any]] = []
    font_violations: list[dict[str, Any]] = []
    for item in TEXT_BOUNDS:
        parent = str(item["parent"])
        if parent and parent in BOX_BOUNDS:
            text_box = (float(item["x"]), float(item["y"]), float(item["w"]), float(item["h"]))
            if not fits_inside(text_box, BOX_BOUNDS[parent]):
                overflow.append(item)
        if int(item["fontSize"]) < LAYOUT.min_font_pt:
            font_violations.append(item)
    return {
        "origin": {"x": LAYOUT.origin_x, "y": LAYOUT.origin_y},
        "slideSize": {"widthIn": LAYOUT.slide_w, "heightIn": LAYOUT.slide_h},
        "exportSize": {"widthPx": LAYOUT.export_w, "heightPx": LAYOUT.export_h},
        "alignment": "left",
        "fontScale": LAYOUT.scale,
        "minFontSizePt": LAYOUT.min_font_pt,
        "trackedTextBoxes": len(TEXT_BOUNDS),
        "overflowCount": len(overflow),
        "fontViolationCount": len(font_violations),
        "overflow": overflow,
        "fontViolations": font_violations,
        "ok": not overflow and not font_violations,
    }


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    export_dir = OUT / "pipeline-overview-export"
    build_deck(PPTX)
    exported = export_with_powerpoint(PPTX, export_dir)
    shutil.copyfile(exported, PNG)
    render = validate_png(PNG)
    layout = validate_layout()
    report = {
        "pptx": str(PPTX.relative_to(ROOT)),
        "png": str(PNG.relative_to(ROOT)),
        "powerPointRawExportPng": str(exported.relative_to(ROOT)),
        "renderValidation": render,
        "layoutValidation": layout,
        "ok": render["hasContent"] and layout["ok"],
    }
    REPORT.write_text(json.dumps(report, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    if not report["ok"]:
        raise SystemExit(json.dumps(report, indent=2, ensure_ascii=False))
    print(json.dumps(report, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
