#!/usr/bin/env python3
from __future__ import annotations

import json
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "artifacts" / "ppt"

SHAPES = [
    {"name": "z00_background_theme", "kind": "rect", "xywh": (0.0, 0.0, 13.333, 7.5), "fill": "F8FAFC", "line": "F8FAFC"},
    {"name": "z10_surface_card", "kind": "roundRect", "xywh": (1.0, 0.9, 9.4, 5.6), "fill": "FFFFFF", "line": "DDE3EA"},
    {"name": "z20_chart_panel", "kind": "rect", "xywh": (1.6, 2.1, 5.7, 3.5), "fill": "E9EEF5", "line": "C7D2DE"},
    {"name": "z30_chart_bar_accent", "kind": "rect", "xywh": (2.2, 3.2, 4.8, 1.0), "fill": "635BFF", "line": "635BFF"},
    {"name": "z40_accent_ribbon", "kind": "rect", "xywh": (6.1, 1.25, 2.9, 4.9), "fill": "00C2A8", "line": "00C2A8"},
    {"name": "z50_callout_overlay", "kind": "roundRect", "xywh": (5.2, 2.15, 4.3, 2.2), "fill": "FFD166", "line": "A97900"},
    {"name": "z60_top_badge", "kind": "roundRect", "xywh": (6.6, 2.55, 2.1, 0.8), "fill": "EF476F", "line": "A51F3C"},
    {"name": "z70_primary_text", "kind": "text", "xywh": (6.85, 2.75, 1.7, 0.35), "text": "TOP", "fill": "FFFFFF", "line": "FFFFFF"},
]

EXPECTED_ORDER = [shape["name"] for shape in SHAPES]
PIXEL_SAMPLES = [
    {"point": (720, 305), "expected": (239, 71, 111), "label": "top badge covers callout/ribbon/chart"},
    {"point": (600, 285), "expected": (255, 209, 102), "label": "callout overlays ribbon/chart"},
    {"point": (660, 500), "expected": (0, 194, 168), "label": "ribbon overlays card"},
    {"point": (310, 365), "expected": (99, 91, 255), "label": "chart bar overlays chart panel"},
]


def rgb(hex_value: str) -> tuple[int, int, int]:
    return tuple(int(hex_value[i:i + 2], 16) for i in (0, 2, 4))


def add_shape(slide, spec: dict) -> None:
    x, y, w, h = spec["xywh"]
    if spec["kind"] == "text":
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        box.name = spec["name"]
        frame = box.text_frame
        frame.clear()
        run = frame.paragraphs[0].add_run()
        run.text = spec["text"]
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*rgb(spec["fill"]))
        return
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if spec["kind"] == "roundRect" else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = spec["name"]
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb(spec["fill"]))
    shape.line.color.rgb = RGBColor(*rgb(spec["line"]))


def create_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for spec in SHAPES:
        add_shape(slide, spec)
    prs.save(path)


def pptx_order(path: Path) -> list[str]:
    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    with zipfile.ZipFile(path) as zf:
        xml = zf.read("ppt/slides/slide1.xml")
    root = ET.fromstring(xml)
    names: list[str] = []
    for node in root.findall(".//p:sp", ns):
        c_nv_pr = node.find(".//p:cNvPr", ns)
        if c_nv_pr is not None:
            name = c_nv_pr.attrib.get("name", "")
            if name.startswith("z"):
                names.append(name)
    return names


def pptx_shape_specs(path: Path) -> list[dict]:
    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }
    with zipfile.ZipFile(path) as zf:
        xml = zf.read("ppt/slides/slide1.xml")
    root = ET.fromstring(xml)
    specs: list[dict] = []
    for node in root.findall(".//p:sp", ns):
        c_nv_pr = node.find(".//p:cNvPr", ns)
        if c_nv_pr is None:
            continue
        name = c_nv_pr.attrib.get("name", "")
        if not name.startswith("z"):
            continue
        off = node.find(".//a:off", ns)
        ext = node.find(".//a:ext", ns)
        srgb = node.find(".//a:solidFill/a:srgbClr", ns)
        prst = node.find(".//a:prstGeom", ns)
        text_nodes = node.findall(".//a:t", ns)
        if off is None or ext is None:
            continue
        emu = 914400
        specs.append({
            "name": name,
            "kind": "text" if text_nodes else ("roundRect" if prst is not None and prst.attrib.get("prst") == "roundRect" else "rect"),
            "xywh": (
                int(off.attrib["x"]) / emu,
                int(off.attrib["y"]) / emu,
                int(ext.attrib["cx"]) / emu,
                int(ext.attrib["cy"]) / emu,
            ),
            "fill": srgb.attrib.get("val", "000000") if srgb is not None else "000000",
            "line": srgb.attrib.get("val", "000000") if srgb is not None else "000000",
            "text": "".join(text.text or "" for text in text_nodes),
        })
    return specs


def to_px(x: float, y: float, w: float, h: float) -> tuple[int, int, int, int]:
    sx, sy = 1280 / 13.333, 720 / 7.5
    return round(x * sx), round(y * sy), round((x + w) * sx), round((y + h) * sy)


def create_visuals(svg_path: Path, png_path: Path, shape_specs: list[dict]) -> None:
    svg_parts = ['<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" viewBox="0 0 1280 720">']
    image = Image.new("RGB", (1280, 720), "white")
    draw = ImageDraw.Draw(image)
    for spec in shape_specs:
        x1, y1, x2, y2 = to_px(*spec["xywh"])
        if spec["kind"] == "text":
            draw.text((x1 + 6, y1 + 4), spec["text"], fill=rgb(spec["fill"]))
            svg_parts.append(f'<text x="{x1 + 6}" y="{y1 + 28}" fill="#{spec["fill"]}" font-size="28" font-weight="700">{spec["text"]}</text>')
            continue
        radius = 18 if spec["kind"] == "roundRect" else 0
        draw.rounded_rectangle((x1, y1, x2, y2), radius=radius, fill=rgb(spec["fill"]), outline=rgb(spec["line"]), width=2)
        svg_parts.append(f'<rect id="{spec["name"]}" x="{x1}" y="{y1}" width="{x2-x1}" height="{y2-y1}" rx="{radius}" fill="#{spec["fill"]}" stroke="#{spec["line"]}" stroke-width="2"/>')
    svg_parts.append("</svg>")
    svg_path.write_text("\n".join(svg_parts), encoding="utf-8")
    image.save(png_path)


def verify_pixels(path: Path) -> list[dict]:
    image = Image.open(path).convert("RGB")
    results = []
    for sample in PIXEL_SAMPLES:
        actual = image.getpixel(sample["point"])
        expected = tuple(sample["expected"])
        ok = sum(abs(a - b) for a, b in zip(actual, expected)) <= 6
        results.append({**sample, "actual": actual, "ok": ok})
    return results


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    pptx_path = OUT / "design_components_z_order_validation.pptx"
    svg_path = OUT / "design_components_z_order_validation.svg"
    png_path = OUT / "design_components_z_order_validation.png"
    report_path = OUT / "z_order_report.json"
    create_pptx(pptx_path)
    order = pptx_order(pptx_path)
    parsed_shapes = pptx_shape_specs(pptx_path)
    create_visuals(svg_path, png_path, parsed_shapes)
    pixel_results = verify_pixels(png_path)
    report = {
        "pptx": str(pptx_path.relative_to(ROOT)),
        "svg": str(svg_path.relative_to(ROOT)),
        "png": str(png_path.relative_to(ROOT)),
        "expectedOrder": EXPECTED_ORDER,
        "actualOrder": order,
        "parsedShapeCount": len(parsed_shapes),
        "parsedShapeNames": [shape["name"] for shape in parsed_shapes],
        "zOrderOk": order == EXPECTED_ORDER,
        "pixelSamples": pixel_results,
        "visualOk": all(item["ok"] for item in pixel_results),
        "note": "PPTX shape order and visual proof are derived from slide1.xml. LibreOffice/PowerPoint rendering is unavailable in this environment, so the PNG is rendered from parsed PPTX shape geometry and colors.",
    }
    report_path.write_text(json.dumps(report, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    if not report["zOrderOk"] or not report["visualOk"]:
        raise SystemExit(json.dumps(report, indent=2, ensure_ascii=False))
    print(json.dumps(report, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
